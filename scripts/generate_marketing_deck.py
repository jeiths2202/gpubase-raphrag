#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenFrame AI KMS 마케팅 전략 기획 PPT 생성기.
TmaxSoft 공식 템플릿 기반, 한국어, 기술 구현 중심 (검증되지 않은 수치 최소화).
codemap.html 캡처 이미지 포함.
"""

import sys
import io
import os
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Emu

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

TEMPLATE_PATH = "TmaxSoft-ppt-template.pptx"
OUTPUT_PATH = "docs/TmaxSoft_AI_KMS_Marketing_Deck.pptx"

# Codemap screenshots (captured by temp/capture_codemap.js)
IMG_GRAPH = "temp/codemap_graph.png"
IMG_SCHEMA = "temp/codemap_schema.png"
IMG_RELATIONS = "temp/codemap_relations.png"

# Architecture diagrams (captured by temp/capture_arch.js)
IMG_ARCH_H = "temp/architecture_horizontal.png"
IMG_ARCH_V = "temp/architecture_vertical.png"


# ---------------------------------------------------------------------------
# Core utilities
# ---------------------------------------------------------------------------

def clone_slide(prs, index):
    src = prs.slides[index]
    new_slide = prs.slides.add_slide(src.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = etree.QName(child.tag).localname
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            sp_tree.remove(child)
    rid_map = {}
    for rel in src.part.rels.values():
        if rel.reltype == RT.SLIDE_LAYOUT:
            for nr in new_slide.part.rels.values():
                if nr.reltype == RT.SLIDE_LAYOUT:
                    rid_map[rel.rId] = nr.rId
                    break
            continue
        try:
            if rel.is_external:
                new_rid = new_slide.part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref
                )
            else:
                new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rid
        except Exception:
            pass
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for child in src.shapes._spTree:
        tag = etree.QName(child.tag).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        nc = deepcopy(child)
        for elem in nc.iter():
            for ak in list(elem.attrib.keys()):
                if R_NS in ak:
                    old = elem.get(ak)
                    if old in rid_map:
                        elem.set(ak, rid_map[old])
        sp_tree.append(nc)
    return new_slide


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def set_text(slide, shape_name, new_text):
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run._r.getparent().remove(run._r)
                    return True
    return False


def add_image(slide, img_path, left, top, width, height):
    """Add an image to a slide if the file exists."""
    if not os.path.exists(img_path):
        return None
    return slide.shapes.add_picture(img_path, left, top, width, height)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

# Image position: right panel Section 3 area (no overlap with text)
# 16:9 aspect ratio, fits within slide bounds
IMG_LEFT = 6000000   # ~6.56in (right panel starts at 6.35in)
IMG_TOP = 4300000    # ~4.70in (below Section 2 ending ~4.26in)
IMG_W = 4500000      # ~4.92in
IMG_H = 2531250      # ~2.77in (16:9 ratio)


def build_cover(prs):
    s = clone_slide(prs, 0)
    set_text(s, "Text 2", "\u00a9 2026. TmaxSoft. All Rights Reserved.")
    set_text(s, "Text 4", "OpenFrame")
    set_text(s, "Text 5", "AI KMS")
    set_text(s, "Text 6", "OpenFrame AI KMS")
    set_text(s, "Text 7", "Platform")
    set_text(s, "Text 8", "\ub808\uac70\uc2dc \ud604\ub300\ud654\ub97c \uc704\ud55c \ub3c4\uba54\uc778 \ud2b9\ud654 AI \uc5d4\uc9c0\ub2c8\uc5b4\ub9c1 \ud50c\ub7ab\ud3fc")
    set_text(s, "Text 9", "Domain-Aware AI Engineering Platform")
    set_text(s, "Text 10", "2026\ub144 3\uc6d4")
    set_text(s, "Text 11", "\uc791\uc131\uc790: \uc2e0\uc774\uc7ac")
    return s


def build_toc(prs):
    s = clone_slide(prs, 1)
    items = {
        "Text 17": "Executive Problem Statement",
        "Text 19": "Why Generic RAG Fails?",
        "Text 21": "Market Timing - Why Now?",
        "Text 23": "Product Overview",
        "Text 25": "Generic RAG vs AI KMS",
        "Text 27": "Core Architecture",
        "Text 29": "4 Core Capabilities",
        "Text 31": "Premium Support & Strategy",
        "Text 33": "Why We Win",
        "Text 35": "Future Roadmap & Messaging",
        "Text 37": "",
        "Text 34": "10",
        "Text 36": "",
    }
    for name, text in items.items():
        set_text(s, name, text)
    return s


def build_content(prs, badge, page, title,
                  left1, left2, left_body,
                  h1, b1, h2, b2, h3="", bul1="", bul2=""):
    s = clone_slide(prs, 2)
    set_text(s, "Text 22", badge)
    set_text(s, "Text 20", title)
    set_text(s, "Text 23", str(page))
    set_text(s, "Text 16", left1)
    set_text(s, "Text 17", left2)
    set_text(s, "Text 18", left_body)
    set_text(s, "Text 6", h1)
    set_text(s, "Text 7", b1)
    set_text(s, "Text 8", h2)
    set_text(s, "Text 9", b2)
    set_text(s, "Text 10", h3)
    set_text(s, "Text 11", bul1)
    set_text(s, "Text 12", bul2)
    set_text(s, "Text 13", "")
    set_text(s, "Text 14", "")
    return s


def build_content_with_image(prs, badge, page, title,
                             left1, left2, left_body,
                             h1, b1, h2, b2, img_path):
    """Content slide with image replacing Section 3."""
    s = build_content(prs, badge, page, title,
                      left1, left2, left_body,
                      h1, b1, h2, b2,
                      "", "", "")  # Clear section 3
    add_image(s, img_path,
              Emu(IMG_LEFT), Emu(IMG_TOP),
              Emu(IMG_W), Emu(IMG_H))
    return s


# Full-slide image (nearly full width, for dedicated architecture slides)
FULL_IMG_W    = 11100000  # ~12.14in (nearly full 13.33in slide width)
FULL_IMG_H    = 6243750   # ~6.83in (16:9 ratio, fits 7.50in slide)
FULL_IMG_LEFT = 546000    # ~0.60in (centered horizontally)
FULL_IMG_TOP  = 550000    # ~0.60in

# Large image covering full right panel (for architecture with text)
BIG_IMG_LEFT = 5650000    # ~6.18in
BIG_IMG_TOP  = 500000     # ~0.55in
BIG_IMG_W    = 6400000    # ~7.00in
BIG_IMG_H    = 3600000    # ~3.94in (16:9 ratio)


def build_full_image_slide(prs, badge, page, title, img_path):
    """Slide with near-full-size image (architecture diagram)."""
    s = build_content(prs, badge, page, title,
                      "", "", "",
                      "", "", "", "",
                      "", "", "")  # Clear all text
    add_image(s, img_path,
              Emu(FULL_IMG_LEFT), Emu(FULL_IMG_TOP),
              Emu(FULL_IMG_W), Emu(FULL_IMG_H))
    return s


def build_content_with_large_image(prs, badge, page, title,
                                   left1, left2, left_body, img_path):
    """Content slide with large image covering entire right panel."""
    s = build_content(prs, badge, page, title,
                      left1, left2, left_body,
                      "", "", "", "",
                      "", "", "")  # Clear all right sections
    add_image(s, img_path,
              Emu(BIG_IMG_LEFT), Emu(BIG_IMG_TOP),
              Emu(BIG_IMG_W), Emu(BIG_IMG_H))
    return s


def build_cards(prs, badge, page, title, cards):
    s = clone_slide(prs, 3)
    set_text(s, "Text 58", badge)
    set_text(s, "Text 56", title)
    set_text(s, "Text 59", str(page))
    t_names = ["Text 6", "Text 11", "Text 16", "Text 21", "Text 26", "Text 31"]
    b_names = ["Text 32", "Text 33", "Text 34", "Text 35", "Text 36", "Text 37"]
    bl_groups = [
        ["Text 38", "Text 39", "Text 40"],
        ["Text 41", "Text 42", "Text 43"],
        ["Text 44", "Text 45", "Text 46"],
        ["Text 47", "Text 48", "Text 49"],
        ["Text 50", "Text 51", "Text 52"],
        ["Text 53", "Text 54", "Text 55"],
    ]
    for i, c in enumerate(cards):
        set_text(s, t_names[i], c["title"])
        set_text(s, b_names[i], c["badge"])
        set_text(s, bl_groups[i][0], c["b1"])
        set_text(s, bl_groups[i][1], c["b2"])
        set_text(s, bl_groups[i][2], c["b3"])
    return s


def build_closing(prs):
    return clone_slide(prs, 4)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation(TEMPLATE_PATH)

    # 1. Cover
    build_cover(prs)
    print("  1  Cover")

    # 2. TOC
    build_toc(prs)
    print("  2  TOC")

    # 3. Executive Problem Statement (HTML slide 3)
    build_content(
        prs, "01", 3, "Executive Problem Statement",
        "\ub808\uac70\uc2dc \uc2dc\uc2a4\ud15c\uc774",
        "\uc9c1\uba74\ud55c '\uad6c\uc870\uc801 \uc704\uae30'",
        "\uba54\uc778\ud504\ub808\uc784/\ub808\uac70\uc2dc \uc804\ubb38\uac00 \uace0\ub839\ud654\ub85c \uc778\ud55c \uae09\uaca9\ud55c \uc740\ud1f4. "
        "\ubb38\uc11c\ud654\ub418\uc9c0 \uc54a\uc740 \uc554\ubb35\uc9c0(Tacit Knowledge)\uc758 \uc18c\uc2e4. "
        "\uc2e0\uaddc \uc778\ub825\uc758 \ub808\uac70\uc2dc \uc5b8\uc5b4(COBOL, ASM) \uae30\ud53c \ud604\uc0c1.",
        "\uc778\ub825 \ubc0f \uc9c0\uc2dd\uc758 \ub2e8\uc808",
        "\ub808\uac70\uc2dc \uc804\ubb38\uac00 \uae09\uaca9\ud55c \uc740\ud1f4\ub85c \ubb38\uc11c\ud654\ub418\uc9c0 \uc54a\uc740 "
        "\uc554\ubb35\uc9c0(Tacit Knowledge)\uac00 \uc18c\uc2e4\ub418\uace0 \uc788\uc2b5\ub2c8\ub2e4. "
        "\uc2e0\uaddc \uc778\ub825\uc740 COBOL, ASM \ub4f1 \ub808\uac70\uc2dc \uc5b8\uc5b4\ub97c \uae30\ud53c\ud558\uc5ec "
        "\uae30\uc220 \uc804\uc218\uac00 \ub2e8\uc808\ub418\ub294 \uad6c\uc870\uc801 \ubb38\uc81c\uac00 \ubc1c\uc0dd\ud569\ub2c8\ub2e4.",
        "\uc6b4\uc601 \ud6a8\uc728\uc131 \uc800\ud558",
        "\uc7a5\uc560 \ubc1c\uc0dd \uc2dc \ub85c\uadf8 \ubd84\uc11d \ubc0f \uc6d0\uc778 \uaddc\uba85\uc5d0 \uacfc\ub2e4 \uc2dc\uac04 \uc18c\uc694. "
        "\ubcf5\uc7a1\ud55c \uc2dc\uc2a4\ud15c \uc758\uc874\uc131\uc73c\ub85c \uc778\ud55c \ubcc0\uacbd \ub9ac\uc2a4\ud06c \uc99d\ub300. "
        "\ub9c8\uc774\uadf8\ub808\uc774\uc158 \ud504\ub85c\uc81d\ud2b8\uc758 \ub192\uc740 \uc2e4\ud328\uc728\uacfc \ube44\uc6a9 \ubd80\ub2f4.",
        "\ud575\uc2ec \uacfc\uc81c",
        "\ub2e8\uc21c \uac80\uc0c9 \ub3c4\uad6c\ub85c\ub294 \ub808\uac70\uc2dc \uc2dc\uc2a4\ud15c\uc758 \uc0c1\ud638\uc758\uc874\uc131 \ud574\uacb0 \ubd88\uac00",
        "\ub3c4\uba54\uc778 \ud2b9\ud654 AI \uae30\ubc18\uc758 \uad6c\uc870\uc801 \uc811\uadfc\uc774 \ud544\uc218",
    )
    print("  3  Executive Problem Statement")

    # 4. Why Generic RAG Fails (HTML slide 4)
    build_content(
        prs, "01", 4, "Why Generic RAG Fails?",
        "\ubc94\uc6a9 AI \ubaa8\ub378\uc774",
        "\uc5d4\ud130\ud504\ub77c\uc774\uc988\uc5d0 \uc2e4\ud328\ud558\ub294 \uc774\uc720",
        "\ubc94\uc6a9 LLM\uc740 OpenFrame \ud2b9\uc720\uc758 \uc5d0\ub7ec \ucf54\ub4dc, "
        "JCL \uc885\uc18d\uc131, \ube44\ud45c\uc900 \ub85c\uadf8 \ud3ec\ub9f7\uc744 "
        "\uc774\ud574\ud558\uc9c0 \ubabb\ud574 \ud53c\uc0c1\uc801\uc778 \ub2f5\ubcc0\ub9cc \uc81c\uacf5\ud569\ub2c8\ub2e4.",
        "1. \ub3c4\uba54\uc778 \ubb38\ub9e5 \ubd80\uc7ac (Context Blindness)",
        "\ubc94\uc6a9 LLM\uc740 OpenFrame \uc5d0\ub7ec \ucf54\ub4dc, JCL \uc885\uc18d\uc131, "
        "\ube44\ud45c\uc900 \ub85c\uadf8 \ud3ec\ub9f7\uc744 \uc774\ud574\ud558\uc9c0 \ubabb\ud574 \ud53c\uc0c1\uc801 \ub2f5\ubcc0\ub9cc \uc81c\uacf5. "
        "\ubbf8\uc158 \ud06c\ub9ac\ud2f0\uceec \uae08\uc735 \uc2dc\uc2a4\ud15c\uc5d0\uc11c \ubd80\uc815\ud655\ud55c \ucf54\ub4dc/\uc124\uc815 \uc81c\uc548\uc740 "
        "\uce58\uba85\uc801 \uc2dc\uc2a4\ud15c \uc7a5\uc560\ub85c \uc9c1\uacb0\ub429\ub2c8\ub2e4.",
        "3. \ube44\uc815\ud615 \ub370\uc774\ud130 \ucc98\ub9ac \ud55c\uacc4 & 4. \ubcf4\uc548 \ubbf8\ube44",
        "\uc218\uc2ed \ub144\uac04 \ucd95\uc801\ub41c \ub364\ud504 \ud30c\uc77c, \ud30c\ud3b8\ud654\ub41c \ub9e4\ub274\uc5bc, "
        "\uad6c\uc870\ud654\ub418\uc9c0 \uc54a\uc740 \ub85c\uadf8\ub97c \uc5f0\uacb0\ud558\uc5ec \ud574\uc11d\ud558\uc9c0 \ubabb\ud569\ub2c8\ub2e4. "
        "\uae30\uc5c5 \ub0b4\ubd80 \ubbfc\uac10 \uc790\uc0b0(\uc18c\uc2a4\ucf54\ub4dc, \ub370\uc774\ud130 \uc2a4\ud0a4\ub9c8)\uc5d0 \ub300\ud55c "
        "\uc815\uad50\ud55c \uc811\uadfc \uc81c\uc5b4(RBAC)\ub3c4 \ubd88\uac00\ub2a5\ud569\ub2c8\ub2e4.",
        "\uacb0\ub860",
        "\ubc94\uc6a9 AI\ub294 \ud45c\uba74\uc801 \ud14d\uc2a4\ud2b8 \ub9e4\uce6d\uc5d0 \uadf8\uce58\ub294 \ud55c\uacc4",
        "\ub3c4\uba54\uc778 \ud2b9\ud654 \ud559\uc2b5 + \uad6c\uc870\uc801 \ud30c\uc2f1\uc774 \ud575\uc2ec \ucc28\ubcc4\uc810",
    )
    print("  4  Why Generic RAG Fails")

    # 5. Market Timing (HTML slide 5)
    build_content(
        prs, "01", 5, "Market Timing - Why Now?",
        "\uc9c0\uae08\uc774 AI \uae30\ubc18",
        "KMS \ub3c4\uc785\uc758 \ucd5c\uc801\uae30\uc778 \uc774\uc720",
        "AI \uc804\ud658 \uac00\uc18d\ud654, \uc804\ubb38 \uc778\ub825 \uc740\ud1f4 \uc784\uacc4\uc810, "
        "\uadf8\ub9ac\uace0 TmaxSoft\uc758 \uae30\uc220 \ucd95\uc801\uc774 \uad50\ucc28\ud558\ub294 "
        "\uc9c0\uae08\uc774 \uacb0\uc815\uc801 \uc2dc\uc810\uc785\ub2c8\ub2e4.",
        "AI \uc804\ud658\uc758 \uac00\uc18d\ud654",
        "\uc804 \uc138\uacc4 \uae30\uc5c5\uc758 AI \ub3c4\uc785\uc774 \uac00\uc18d\ud654\ub418\uba70, "
        "\ub808\uac70\uc2dc \ud604\ub300\ud654 \uc2dc\uc7a5\uc5d0\uc11c\ub3c4 \uc790\ub3d9\ud654 \uc694\uad6c\uac00 \ud3ed\ubc1c\uc801\uc73c\ub85c "
        "\uc99d\uac00\ud558\uace0 \uc788\uc2b5\ub2c8\ub2e4. \uae30\uc5c5\ub4e4\uc740 \ub2e8\uc21c \ubbf8\ub4e4\uc6e8\uc5b4\ub97c \ub118\uc5b4 "
        "AI \uae30\ubc18 \uc778\ud154\ub9ac\uc804\ud2b8 \ud50c\ub7ab\ud3fc\uc744 \uc694\uad6c\ud569\ub2c8\ub2e4.",
        "\uc9c0\uc2dd \uc99d\ubc1c \uc784\uacc4\uc810",
        "\ubca0\uc774\ube44\ubd80\uba38 \uc5d4\uc9c0\ub2c8\uc5b4\uc758 \ub300\uac70 \uc740\ud1f4\uac00 \uc2dc\uc791\ub418\uc5c8\uc2b5\ub2c8\ub2e4. "
        "\uc9c0\uae08 \uc9c0\uc2dd\uc744 \uc790\uc0b0\ud654\ud558\uc9c0 \uc54a\uc73c\uba74 \uc601\uad6c\uc801\uc778 \uae30\uc220 \ubd80\ucc44\uac00 \ub429\ub2c8\ub2e4. "
        "OpenFrame \uc6b4\uc601 \uacbd\ud5d8\uacfc \ud2b8\ub7ec\ube14\uc288\ud305 \ub178\ud558\uc6b0\ub97c AI\uc5d0 "
        "\ud559\uc2b5\uc2dc\ud0a4\ub294 \uac83\uc774 \uc2dc\uae09\ud569\ub2c8\ub2e4.",
        "TmaxSoft\uc758 \ub3c4\uc57d",
        "\ub2e8\uc21c \ubbf8\ub4e4\uc6e8\uc5b4 \ubca4\ub354 \u2192 AI \uae30\uc220 \uacb0\ud569 \uc778\ud154\ub9ac\uc804\ud2b8 \ud50c\ub7ab\ud3fc \uae30\uc5c5",
        "\uc218\uc2ed \ub144 \ucd95\uc801\ub41c \ub3c4\uba54\uc778 \ub370\uc774\ud130\uac00 \ucd5c\uac15\uc758 \uacbd\uc7c1 \uc6b0\uc704",
    )
    print("  5  Market Timing")

    # 6. Product Overview + codemap schema (HTML slide 6)
    build_content_with_image(
        prs, "02", 6, "Product Overview",
        "Domain-Aware AI",
        "Engineering Platform",
        "\ub2e8\uc21c \ucc57\ubd07\uc774 \uc544\ub2cc, OpenFrame \ud2b9\ud654 "
        "\uc5d4\uc9c0\ub2c8\uc5b4\ub9c1 \ud50c\ub7ab\ud3fc. \ub9e4\ub274\uc5bc, \uc18c\uc2a4\ucf54\ub4dc \ud328\ud134, "
        "\ud2b8\ub7ec\ube14\uc288\ud305 \uac00\uc774\ub4dc\ub97c \ud559\uc2b5\ud558\uc5ec "
        "\uc5d4\uc9c0\ub2c8\uc5b4 \uc218\uc900\uc758 \ub2f5\ubcc0\uc744 \uc81c\uacf5\ud569\ub2c8\ub2e4.",
        "RAFT \uae30\ubc18 \ub3c4\uba54\uc778 \ud2b9\ud654 \ud559\uc2b5",
        "\ub9e4\ub274\uc5bc, \ud2b8\ub7ec\ube14\uc288\ud305 \uac00\uc774\ub4dc, \uc18c\uc2a4\ucf54\ub4dc \ud328\ud134\uc744 \ud559\uc2b5\ud558\uc5ec "
        "\uc5d4\uc9c0\ub2c8\uc5b4 \uc218\uc900 \ub2f5\ubcc0\uc744 \uc81c\uacf5\ud569\ub2c8\ub2e4. "
        "QLoRA 3-Phase (CPT\u2192SFT\u2192DPO) \ud30c\uc774\ud504\ub77c\uc778\uc73c\ub85c "
        "\uc81c\ud488\ubcc4 \uc804\ubb38 AI \uc5b4\ub311\ud130\ub97c \uad6c\ucd95\ud569\ub2c8\ub2e4.",
        "\ud558\uc774\ube0c\ub9ac\ub4dc \ucd94\ub860 \uc5d4\uc9c4 & \ub0b4\uc7a5\ud615 \ud30c\uc11c",
        "\ud655\ub960\uc801 LLM \ub2f5\ubcc0\uacfc \uacb0\uc815\ub860\uc801(Rule-based) \uc2dc\uc2a4\ud15c \ubd84\uc11d\uc744 "
        "\uacb0\ud569\ud558\uc5ec \uc815\ud655\ub3c4\ub97c \uadf9\ub300\ud654\ud569\ub2c8\ub2e4. "
        "ASM, COBOL, JCL \ub4f1 \ub808\uac70\uc2dc \uc5b8\uc5b4 \ud30c\uc11c\ub97c \ub0b4\uc7a5\ud558\uc5ec "
        "\ucf54\ub4dc \uc5f0\uad00 \uad00\uacc4\ub97c \uad6c\uc870\uc801\uc73c\ub85c \ubd84\uc11d\ud569\ub2c8\ub2e4.",
        IMG_SCHEMA,
    )
    print("  6  Product Overview + codemap schema")

    # 7. Generic RAG vs AI KMS (HTML slide 7)
    build_cards(
        prs, "02", 7, "Generic RAG vs TmaxSoft AI KMS",
        [
            {"title": "\uc9c0\uc2dd \uc774\ud574\ub3c4", "badge": "\uc2ec\uce35",
             "b1": "Generic: \ud45c\uba74\uc801 \ud14d\uc2a4\ud2b8 \ub9e4\uce6d (Keyword/Vector)",
             "b2": "AI KMS: \uc2ec\uce35\uc801 \uad6c\uc870 \uc774\ud574 (Code & Log)",
             "b3": "\uc5d4\uc9c4 \ub808\ubca8 \uc5d0\ub7ec \ucf54\ub4dc/JCL \uc885\uc18d\uc131 \ud30c\uc545"},
            {"title": "\ub370\uc774\ud130 \ucc98\ub9ac", "badge": "Parser",
             "b1": "Generic: \uc77c\ubc18 \ubb38\uc11c (PDF, Word) \uc911\uc2ec",
             "b2": "AI KMS: OpenFrame \ub85c\uadf8, \ub364\ud504, \uc18c\uc2a4\ucf54\ub4dc \ud30c\uc2f1",
             "b3": "\ub0b4\uc7a5\ud615 \ud30c\uc11c\ub85c \ub808\uac70\uc2dc \uc5b8\uc5b4 \uad6c\uc870 \ubd84\uc11d"},
            {"title": "\uc2e0\ub8b0\uc131", "badge": "RAFT",
             "b1": "Generic: \ud658\uac01(Hallucination) \uac00\ub2a5\uc131 \ub192\uc74c",
             "b2": "AI KMS: RAFT & \uadfc\uac70 \uae30\ubc18 \uac80\uc99d\uc73c\ub85c \ucd5c\uc18c\ud654",
             "b3": "E2E Hallucination \ud14c\uc2a4\ud2b8 \uc790\ub3d9\ud654 \uac80\uc99d"},
            {"title": "\uc8fc \ud65c\uc6a9\ucc98", "badge": "Engine",
             "b1": "Generic: \uc77c\ubc18 \uc5c5\ubb34 \uc9c0\uc6d0, \uaddc\uc815 \uac80\uc0c9",
             "b2": "AI KMS: \uc7a5\uc560 \ubd84\uc11d, \ub9c8\uc774\uadf8\ub808\uc774\uc158, \uc131\ub2a5 \ud29c\ub2dd",
             "b3": "Multi-Agent\uac00 \ubcf5\ud569 \uc5c5\ubb34\ub97c \uc790\uc728 \uc218\ud589"},
            {"title": "\ubcf4\uc548 \ubaa8\ub378", "badge": "RBAC",
             "b1": "Generic: \ubc94\uc6a9 \ubcf4\uc548 \uc815\ucc45",
             "b2": "AI KMS: \uc5d4\ud130\ud504\ub77c\uc774\uc988\uae09 RBAC & On-Premise",
             "b3": "SSO/OAuth + \uac10\uc0ac \ub85c\uadf8 + \ub2e4\uc911 \uc778\uc99d"},
            {"title": "\ud559\uc2b5 \ubc29\uc2dd", "badge": "QLoRA",
             "b1": "Generic: \ubc94\uc6a9 \uc0ac\uc804\ud559\uc2b5 \ubaa8\ub378 \uadf8\ub300\ub85c \uc0ac\uc6a9",
             "b2": "AI KMS: 3-Phase QLoRA \ub3c4\uba54\uc778 \ud2b9\ud654 \ud559\uc2b5",
             "b3": "\uc81c\ud488\ubcc4 \uc5b4\ub311\ud130\ub85c \uc815\ud655\ud55c \uc804\ubb38 \uc9c0\uc2dd \ub2f5\ubcc0"},
        ],
    )
    print("  7  Generic RAG vs AI KMS (6 cards)")

    # 8. Core Architecture + codemap graph (HTML slide 8)
    build_content_with_image(
        prs, "02", 8, "Core Architecture",
        "AI \ud30c\uc774\ud504\ub77c\uc778",
        "\uc544\ud0a4\ud14d\ucc98",
        "Input \u2192 Parser & Embedding \u2192 Hybrid Search "
        "\u2192 Reasoning \u2192 Verified Output. "
        "\uc2e0\ub8b0\ud560 \uc218 \uc788\ub294 AI \ud30c\uc774\ud504\ub77c\uc778\uc744 \uad6c\ud604\ud569\ub2c8\ub2e4.",
        "Intelligent Brain",
        "RAFT \uae30\ubc18 \ub3c4\uba54\uc778 \ud2b9\ud654 \ud559\uc2b5\uc73c\ub85c \uc9c8\uc758 \uc758\ub3c4\ub97c \ud30c\uc545\ud558\uace0, "
        "QueryRouter\uac00 LLM \uc5c6\uc774 \uc81c\ud488\uc744 \uc790\ub3d9 \ub77c\uc6b0\ud305\ud569\ub2c8\ub2e4. "
        "ProductContextMemory\ub85c \uc138\uc158 \uac04 \ucee8\ud14d\uc2a4\ud2b8\ub97c \uc601\uc18d\ud654\ud558\uc5ec "
        "\uc5f0\uc18d\uc801\uc778 \ub300\ud654 \ud750\ub984\uc744 \uc9c0\uc6d0\ud569\ub2c8\ub2e4.",
        "Analysis Engine",
        "Legacy Parser: COBOL/JCL/ASM \uad6c\uc870 \ubd84\uc11d. "
        "Hybrid Search: Vector + Graph + BM25 3\uc911 \uac80\uc0c9 \uacb0\ud569. "
        "Neo4j Graph DB: Chunk-Entity \uad00\uacc4 \ubaa8\ub378\ub9c1\uc73c\ub85c "
        "\uc5f0\uad00 \uad00\uacc4\ub97c \uc2dc\uac01\ud654\ud558\uace0 \ud0d0\uc0c9\ud569\ub2c8\ub2e4.",
        IMG_GRAPH,
    )
    print("  8  Core Architecture + codemap graph")

    # 9. Enterprise Architecture (horizontal component view)
    build_content_with_large_image(
        prs, "02", 9, "Enterprise Architecture",
        "컴포넌트 기반",
        "플랫폼 전체 구조도",
        "Legacy Mainframe에서 OpenFrame Integration을 거쳐 "
        "AI KMS Platform까지의 전체 데이터/제어 흐름을 "
        "컴포넌트 단위로 시각화합니다. "
        "각 요소의 내부 기술 스택과 상호 연결을 확인할 수 있습니다.",
        IMG_ARCH_H,
    )
    print("  9  Enterprise Architecture (horizontal)")

    # 10. 4 Core Capabilities (HTML slide 11)
    build_cards(
        prs, "03", 10, "4 Core Capabilities",
        [
            {"title": "AI Technical Support", "badge": "Support",
             "b1": "\uc790\uc5f0\uc5b4 \uc9c8\uc758\uc751\ub2f5, \uc5d0\ub7ec \ucf54\ub4dc \ud574\uacb0",
             "b2": "\ub9e4\ub274\uc5bc \uae30\ubc18 \uac00\uc774\ub4dc + \uadfc\uac70 \ub9c1\ud06c \uc81c\uacf5",
             "b3": "L1 \uae30\uc220\uc9c0\uc6d0 \uc790\ub3d9\ud654\ub85c \ub300\uae30 \uc2dc\uac04 \ud574\uc18c"},
            {"title": "AI Log & Core Analysis", "badge": "Analysis",
             "b1": "JOB \ub85c\uadf8 \uad6c\uc870 \ud30c\uc2f1 \ubc0f \uc5d4\uc9c4 \ub808\ubca8 \ubd84\uc11d",
             "b2": "Core Dump \uc2a4\ud0dd \ud2b8\ub808\uc774\uc2a4 \uc790\ub3d9 \ubd84\uc11d",
             "b3": "\uadfc\ubcf8 \uc6d0\uc778 \ub3c4\ucd9c \ubc0f \uc870\uce58 \uac00\uc774\ub4dc \uc81c\uacf5"},
            {"title": "AI Migration Design", "badge": "Migrate",
             "b1": "\uc18c\uc2a4\ucf54\ub4dc Call Graph \uc790\ub3d9 \uc0dd\uc131",
             "b2": "\ube44\ud638\ud658 \ucf54\ub4dc \ud0d0\uc9c0 \ubc0f \ubcc0\ud658 \uac00\uc774\ub4dc",
             "b3": "Dead Code \uc2dd\ubcc4\ub85c \ub9c8\uc774\uadf8\ub808\uc774\uc158 \ucd5c\uc801\ud654"},
            {"title": "AI Asset Documentation", "badge": "Docs",
             "b1": "COBOL/PL/I \uc5ed\ubd84\uc11d \u2192 \ube44\uc988\ub2c8\uc2a4 \ub85c\uc9c1 \uc124\uba85",
             "b2": "\ub370\uc774\ud130 \ud750\ub984\ub3c4(Flowchart) & ERD \uc0dd\uc131",
             "b3": "\ucf54\ub4dc \ubcc0\uacbd \uc2dc \ubb38\uc11c \uc790\ub3d9 \uc5c5\ub370\uc774\ud2b8"},
            {"title": "Premium Support", "badge": "Premium",
             "b1": "AI Instant Solve: \uc774\uc288 \ub4f1\ub85d \uc2dc 1\ucc28 \ud574\uacb0",
             "b2": "AI-Assisted Expert: \ubbf8\ud574\uacb0 \uac74 \uc804\ubb38\uac00 \uc774\uad00",
             "b3": "\uc6d0\uaca9 \ud611\uc5c5 + \uace0\uac1d \uc804\uc6a9 AI \ubaa8\ub378 \ud29c\ub2dd"},
            {"title": "Enterprise Security", "badge": "Security",
             "b1": "\uc5ed\ud560 \uae30\ubc18 \uc811\uadfc \uc81c\uc5b4 (RBAC)",
             "b2": "\uac10\uc0ac \ub85c\uadf8 \ubc0f \ub2e4\uc911 \uc778\uc99d (MFA)",
             "b3": "SSO/OAuth \ud1b5\ud569, On-Premise \ubc30\ud3ec"},
        ],
    )
    print(" 10  4 Core Capabilities (6 cards)")

    # 11. AI Technical Support (HTML slide 12)
    build_content(
        prs, "03", 11, "AI Technical Support",
        "24/7/365 \uac00\ub3d9\ub418\ub294",
        "\uac00\uc0c1\uc758 \uc218\uc11d \uc5d4\uc9c0\ub2c8\uc5b4",
        "\ub2e8\uc21c \ubb38\uc758\uc5d0 \uace0\uae09 \uc5d4\uc9c0\ub2c8\uc5b4 \ub9ac\uc18c\uc2a4 \ub0ad\ube44, "
        "\uc57c\uac04/\ud734\uc77c \uc9c0\uc6d0 \uacf5\ubc31. "
        "AI\uac00 \uc790\uc5f0\uc5b4 \uc9c8\uc758\uc5d0 \uc989\uc2dc \ub2f5\ubcc0\ud558\uace0 "
        "\uadfc\uac70 \ub9c1\ud06c\ub97c \ud568\uaed8 \uc81c\uacf5\ud569\ub2c8\ub2e4.",
        "Problem",
        "\ub2e8\uc21c \ubb38\uc758\uc5d0 \uace0\uae09 \uc5d4\uc9c0\ub2c8\uc5b4 \ub9ac\uc18c\uc2a4 \ub0ad\ube44. "
        "\uc57c\uac04/\ud734\uc77c \uc9c0\uc6d0 \uacf5\ubc31 \ubc1c\uc0dd. "
        "IMS(\uc774\uc288 \uad00\ub9ac \uc2dc\uc2a4\ud15c) \ub370\uc774\ud130 2\ub9cc \uac74 \uc774\uc0c1 \ud65c\uc6a9 \ubd88\uac00.",
        "Solution",
        "'OSC1234 \uc5d0\ub7ec \uc870\uce58\ubc95 \uc54c\ub824\uc918' \ub4f1 \uc790\uc5f0\uc5b4 \uc9c8\uc758 \uc989\uc2dc \ub2f5\ubcc0. "
        "\ub2f5\ubcc0 \uadfc\uac70(\ub9e4\ub274\uc5bc \ud398\uc774\uc9c0, \uc720\uc0ac Ticket) \ub9c1\ud06c \uc81c\uacf5. "
        "Agentic RAG Two-Stage \uac80\uc0c9\uc73c\ub85c \uc815\ud655\ub3c4 \uadf9\ub300\ud654.",
        "Impact",
        "L1 \uae30\uc220\uc9c0\uc6d0 \uc790\ub3d9\ud654, \uace0\uac1d \ub300\uae30 \uc2dc\uac04 \ucd5c\uc18c\ud654",
        "\uc5d4\uc9c0\ub2c8\uc5b4 \ub9ac\uc18c\uc2a4\ub97c \uace0\ub09c\ub3c4 \uc774\uc288\uc5d0 \uc9d1\uc911",
    )
    print(" 11  AI Technical Support")

    # 12. AI Log & Core Analysis (HTML slide 13)
    build_content(
        prs, "03", 12, "AI Log & Core Analysis",
        "OpenFrame \uc804\ubb38\uac00",
        "Agent\uac00 \uc5d4\uc9c4 \ub808\ubca8 \uc815\ubc00 \ubd84\uc11d",
        "\uc218\ucc9c \uc904\uc758 JOB \ub85c\uadf8\ub97c \uc0ac\ub78c\uc774 \uc9c1\uc811 \ubd84\uc11d, "
        "\uc5d0\ub7ec \uc6d0\uc778 \ud30c\uc545\uc5d0 \uc218 \uc2dc\uac04 \uc18c\uc694. "
        "AI Agent\uac00 \ub85c\uadf8\ub97c \uad6c\uc870\uc801\uc73c\ub85c \ud30c\uc2f1\ud558\uc5ec "
        "\uadfc\ubcf8 \uc6d0\uc778\uc744 \ub3c4\ucd9c\ud569\ub2c8\ub2e4.",
        "Problem",
        "\uc218\ucc9c \uc904\uc758 JOB \ub85c\uadf8\ub97c \uc0ac\ub78c\uc774 \uc9c1\uc811 \ubd84\uc11d. "
        "\uc5d0\ub7ec \uc6d0\uc778 \ud30c\uc545\uc5d0 \uc218 \uc2dc\uac04 \uc18c\uc694. "
        "\ubcf5\uc7a1\ud55c \uc2dc\uc2a4\ud15c \uc758\uc874\uc131\uc73c\ub85c \uc624\uc9c4 \ube48\ubc88.",
        "Solution",
        "OpenFrame \uc804\ubb38\uac00 AI Agent\uac00 JOB \ub85c\uadf8\ub97c \uad6c\uc870\uc801\uc73c\ub85c \ud30c\uc2f1 \ubc0f \ubd84\uc11d. "
        "\uc5d0\ub7ec \ucf54\ub4dc\uc640 \ucee8\ud14d\uc2a4\ud2b8\ub97c \uc5d4\uc9c4 \ub808\ubca8\uc5d0\uc11c \uc815\ud655\ud788 \ud310\ub2e8\ud558\uc5ec "
        "\uadfc\ubcf8 \uc6d0\uc778\uc744 \ub3c4\ucd9c. Core Dump \ubc1c\uc0dd \uc2dc "
        "\uc2a4\ud0dd \ud2b8\ub808\uc774\uc2a4 \uc790\ub3d9 \ubd84\uc11d \ubc0f \ubb38\uc81c \ubaa8\ub4c8 \uc2dd\ubcc4.",
        "Impact",
        "JOB \uc2e4\ud328 \uc6d0\uc778 \ubd84\uc11d \uc2dc\uac04 \ub300\ud3ed \ub2e8\ucd95",
        "\uc5d4\uc9c4 \ub808\ubca8 \uc815\ud655\ub3c4 \ud5a5\uc0c1",
    )
    print(" 12  AI Log & Core Analysis")

    # 13. AI Migration Design + codemap relations (HTML slide 14)
    build_content_with_image(
        prs, "03", 13, "AI Migration Design",
        "\uc131\uacf5\uc801\uc778 OpenFrame",
        "\uc804\ud658\uc744 \uc704\ud55c AI \ub124\ube44\uac8c\uc774\ud130",
        "\ubc29\ub300\ud55c AS-IS \uc790\uc0b0 \ud30c\uc545 \ub09c\ud56d, "
        "\ube44\ud638\ud658 \uc694\uc18c \uc0ac\uc804 \uc2dd\ubcc4 \ub204\ub77d. "
        "\uc18c\uc2a4\ucf54\ub4dc \uc2a4\uce94\uc73c\ub85c Call Graph \uc0dd\uc131 \ubc0f "
        "Dead Code \uc2dd\ubcc4\ub85c \ub9c8\uc774\uadf8\ub808\uc774\uc158 \ucd5c\uc801\ud654.",
        "Problem & Solution",
        "\ubc29\ub300\ud55c AS-IS \uc790\uc0b0 \ud30c\uc545 \ub09c\ud56d, \ube44\ud638\ud658 \uc694\uc18c \uc0ac\uc804 \uc2dd\ubcc4 \ub204\ub77d. "
        "\uc18c\uc2a4\ucf54\ub4dc \uc2a4\uce94\uc73c\ub85c \ud504\ub85c\uadf8\ub7a8 \uac04 \ud638\ucd9c \uad00\uacc4(Call Graph) \uc0dd\uc131. "
        "\uc804\ud658 \uc2dc \uc218\uc815 \ud544\uc694\ud55c \ucf54\ub4dc \ube14\ub85d \uc790\ub3d9 \uc2dd\ubcc4 \ubc0f \ubcc0\ud658 \uac00\uc774\ub4dc \uc81c\uacf5. "
        "\ubbf8\uc0ac\uc6a9 \uc790\uc0b0(Dead Code) \uc2dd\ubcc4\ub85c \ub9c8\uc774\uadf8\ub808\uc774\uc158 \ub300\uc0c1 \ucd5c\uc801\ud654.",
        "Impact",
        "\ucd08\uae30 \ubd84\uc11d \uae30\uac04 \ub2e8\ucd95, \uc804\ud658 \ub9ac\uc2a4\ud06c \uc0ac\uc804 \uc81c\uac70. "
        "\ub808\uac70\uc2dc \uc2dc\uc2a4\ud15c\uc758 \ud22c\uba85\uc131 \ud655\ubcf4.",
        IMG_RELATIONS,
    )
    print(" 13  AI Migration Design + codemap relations")

    # 14. AI Asset Documentation (HTML slide 15)
    build_content(
        prs, "03", 14, "AI Asset Documentation",
        "\ube14\ub799\ubc15\uc2a4\ud654\ub41c",
        "\ub808\uac70\uc2dc \uc2dc\uc2a4\ud15c\uc758 \ud22c\uba85\uc131 \ud655\ubcf4",
        "\ubb38\uc11c\uc640 \uc2e4\uc81c \ucf54\ub4dc \ubd88\uc77c\uce58, "
        "\uac1c\ubc1c\uc790 \ud1f4\uc0ac\ub85c \ube44\uc988\ub2c8\uc2a4 \ub85c\uc9c1 \uc720\uc2e4. "
        "AI\uac00 \ucf54\ub4dc\ub97c \uc5ed\ubd84\uc11d\ud558\uc5ec \ube44\uc988\ub2c8\uc2a4 \ub85c\uc9c1 "
        "\uc124\uba85\uacfc \ubb38\uc11c\ub97c \uc790\ub3d9 \uc0dd\uc131\ud569\ub2c8\ub2e4.",
        "Problem",
        "\ubb38\uc11c\uc640 \uc2e4\uc81c \ucf54\ub4dc \ubd88\uc77c\uce58. "
        "\uac1c\ubc1c\uc790 \ud1f4\uc0ac\ub85c \ube44\uc988\ub2c8\uc2a4 \ub85c\uc9c1 \uc720\uc2e4. "
        "\uc2e0\uaddc \uc778\ub825\uc758 \uc5c5\ubb34 \uc801\uc751\uc5d0 \uc7a5\uae30\uac04 \uc18c\uc694.",
        "Solution",
        "COBOL/PL/I \ucf54\ub4dc \uc5ed\ubd84\uc11d\ud558\uc5ec \ube44\uc988\ub2c8\uc2a4 \ub85c\uc9c1 \uc124\uba85(\ud55c/\uc601) \uc0dd\uc131. "
        "\ub370\uc774\ud130 \ud750\ub984\ub3c4(Flowchart) \ubc0f ERD \uc790\ub3d9 \uc0dd\uc131. "
        "\ucf54\ub4dc \ubcc0\uacbd \uc2dc \ubb38\uc11c \uc790\ub3d9 \uc5c5\ub370\uc774\ud2b8 (Live Docs)\ub85c "
        "\ubb38\uc11c\uc640 \uc2e4\uc81c \ucf54\ub4dc \ubd88\uc77c\uce58 \ubb38\uc81c\ub97c \ud574\uacb0\ud569\ub2c8\ub2e4.",
        "Impact",
        "\uc720\uc9c0\ubcf4\uc218 \uc0dd\uc0b0\uc131 \ud5a5\uc0c1, \uc2e0\uaddc \uc778\ub825\uc758 \ube60\ub978 \uc5c5\ubb34 \uc801\uc751",
        "\ube14\ub799\ubc15\uc2a4\ud654\ub41c \ucf54\ub4dc\uc758 \ud22c\uba85\uc131 \ud655\ubcf4",
    )
    print(" 14  AI Asset Documentation")

    # 15. Premium Support Model (HTML slide 16)
    build_content(
        prs, "04", 15, "Premium Support Model",
        "AI-First,",
        "Human-Expert Second",
        "\uace0\uac1d \uc774\uc288 \ub4f1\ub85d \uc2dc AI\uac00 \uc989\uc2dc \ub85c\uadf8 \ubd84\uc11d \ud6c4 "
        "1\ucc28 \ud574\uacb0\ucc45 \ubc0f \uadfc\uac70 \ubb38\uc11c\ub97c \uc81c\uc2dc\ud569\ub2c8\ub2e4. "
        "\ubbf8\ud574\uacb0 \ubcf5\uc7a1 \uc774\uc288\ub294 AI \ubd84\uc11d \ub9ac\ud3ec\ud2b8\uc640 \ud568\uaed8 "
        "\uc804\ubb38 \uc5d4\uc9c0\ub2c8\uc5b4\uc5d0\uac8c \uc774\uad00\ub429\ub2c8\ub2e4.",
        "Step 1: AI Instant Solve",
        "\uace0\uac1d \uc774\uc288 \ub4f1\ub85d \uc2dc AI\uac00 \uc989\uc2dc \ub85c\uadf8 \ubd84\uc11d \ud6c4 "
        "1\ucc28 \ud574\uacb0\ucc45 \ubc0f \uadfc\uac70 \ubb38\uc11c\ub97c \uc81c\uc2dc\ud569\ub2c8\ub2e4.",
        "Step 2: AI-Assisted Expert",
        "\ubbf8\ud574\uacb0 \ubcf5\uc7a1 \uc774\uc288\ub294 AI \ubd84\uc11d \ub9ac\ud3ec\ud2b8\uc640 \ud568\uaed8 "
        "\uc804\ubb38 \uc5d4\uc9c0\ub2c8\uc5b4\uc5d0\uac8c \uc774\uad00\ub418\uc5b4 \uc815\ubc00 \ubd84\uc11d.",
        "Enterprise Value",
        "\uc6d0\uaca9 \ud611\uc5c5 \ud50c\ub7ab\ud3fc(\ud654\uba74 \uacf5\uc720/\ub77c\uc774\ube0c \ubd84\uc11d) \uc81c\uacf5",
        "\uace0\uac1d \uc804\uc6a9 AI \ubaa8\ub378 \ud29c\ub2dd\uc73c\ub85c \ub77c\uc774\uc120\uc2a4 \uac00\uce58 \uc81c\uace0",
    )
    print(" 15  Premium Support Model")

    # 16. Strategic Positioning + vertical architecture (HTML slide 17)
    build_content_with_image(
        prs, "04", 16, "Strategic Positioning",
        "OpenFrame AI",
        "Operating Layer",
        "단순한 지원 도구(Tool)가 아닌, "
        "OpenFrame 가치를 극대화하는 "
        "'운영 레이어(Operating Layer)'로 "
        "포지셔닝합니다.",
        "3층 아키텍처",
        "Customer Applications (상층) "
        "→ OpenFrame AI KMS Platform - Intelligence Layer (중층) "
        "→ OpenFrame Core System - Runtime (하층). "
        "AI KMS는 단순 Tool이 아닌 Operating Layer로 포지셔닝.",
        "Lock-in & 가치 제고",
        "레거시 현대화의 필수 패키지로 제안하여 "
        "라이선스 가치 제고 및 Lock-in 효과 강화.",
        IMG_ARCH_V,
    )
    print(" 16  Strategic Positioning + vertical arch")

    # 17. Corporate Transformation (HTML slide 18)
    build_content(
        prs, "04", 17, "Corporate Transformation",
        "Product Vendor \u2192",
        "AI Platform Company",
        "\uc81c\ud488 \ub77c\uc774\uc120\uc2a4 \uc911\uc2ec \ub9e4\ucd9c\uc5d0\uc11c "
        "AI \uc11c\ube44\uc2a4 \uad6c\ub3c5(SaaS) \ub9e4\ucd9c\ub85c \uc804\ud658. "
        "\uc778\ub825 \uc758\uc874\uc801 \uc9c0\uc6d0\uc5d0\uc11c "
        "\ub370\uc774\ud130/AI \uae30\ubc18 \uc9c0\uc6d0\uc73c\ub85c \ud601\uc2e0.",
        "Current State: Middleware Product Vendor",
        "\uc81c\ud488 \ub77c\uc774\uc120\uc2a4 \uc911\uc2ec \ub9e4\ucd9c. "
        "\uc778\ub825 \uc758\uc874\uc801 \uae30\uc220 \uc9c0\uc6d0. "
        "\ubc18\uc751\ud615(Reactive) \uc7a5\uc560 \ub300\uc751.",
        "Future State: AI Modernization Platform",
        "AI \uc11c\ube44\uc2a4 \uad6c\ub3c5(SaaS) \ub9e4\ucd9c \ucd94\uac00. "
        "\ub370\uc774\ud130/AI \uae30\ubc18 \uae30\uc220 \uc9c0\uc6d0. "
        "\uc608\uce21\ud615(Proactive) \uc2dc\uc2a4\ud15c \uad00\ub9ac \uc2e4\ud604.",
        "Transformation Value",
        "\uae00\ub85c\ubc8c \uacbd\uc7c1\ub825 \ud655\ubcf4\uc640 \uc9c0\uc18d \uc131\uc7a5 \uae30\ubc18 \ub9c8\ub828",
        "AI \uae30\ubc18 \ub808\uac70\uc2dc \ud604\ub300\ud654\uc758 \uae00\ub85c\ubc8c \ub9ac\ub354",
    )
    print(" 17  Corporate Transformation")

    # 18. Why We Win (HTML slide 19)
    build_cards(
        prs, "04", 18, "Why We Win",
        [
            {"title": "\ub3c5\ubcf4\uc801\uc778 \ub3c4\uba54\uc778 \ub370\uc774\ud130", "badge": "Data",
             "b1": "\uc218\uc2ed \ub144\uac04 \ucd95\uc801\ub41c OpenFrame \uae30\uc220 \ubb38\uc11c",
             "b2": "\uc774\uc288 \ub370\uc774\ud130, \ub9c8\uc774\uadf8\ub808\uc774\uc158 \ub178\ud558\uc6b0",
             "b3": "\ubcf5\uc81c \ubd88\uac00\ub2a5\ud55c \ub3c5\uc810\uc801 \ud559\uc2b5 \uc790\uc0b0"},
            {"title": "Deep Integration", "badge": "Parser",
             "b1": "OpenFrame \uad6c\uc870 \uc644\ubcbd \uc774\ud574 \ud30c\uc11c \ub0b4\uc7a5",
             "b2": "\uac89\ud56b\uae30\uc2dd RAG \uc194\ub8e8\uc158\uacfc \ucc28\uc6d0\uc774 \ub2e4\ub978 \ubd84\uc11d",
             "b3": "ASM/COBOL/JCL \ucf54\ub4dc \uad6c\uc870\uc801 \ud574\uc11d"},
            {"title": "\uad6c\uc870\uc801 \ucd94\ub860 (Reasoning)", "badge": "Graph",
             "b1": "\ud14d\uc2a4\ud2b8 \uac80\uc0c9\uc744 \ub118\uc5b4 \uc2e4\ud589 \ud750\ub984 \uc774\ud574",
             "b2": "\uc758\uc874\uc131 \uadf8\ub798\ud504 \uae30\ubc18 \ub17c\ub9ac\uc801 \ucd94\ub860",
             "b3": "Neo4j Chunk-Entity \uad00\uacc4 \ubaa8\ub378\ub9c1"},
            {"title": "\uc5d4\uc9c0\ub2c8\uc5b4\ub9c1 \uc790\ub3d9\ud654", "badge": "Action",
             "b1": "\ub2e8\uc21c \uac80\uc0c9\uc744 \ub118\uc5b4 \uc2e4\uc81c \ucf54\ub4dc \ubcc0\ud658",
             "b2": "\uc124\uc815 \ubcc0\uacbd, \ub9c8\uc774\uadf8\ub808\uc774\uc158 \uac00\uc774\ub4dc \uc790\ub3d9 \uc0dd\uc131",
             "b3": "Multi-Agent \uae30\ubc18 \ubcf5\ud569 \uc5c5\ubb34 \uc218\ud589"},
            {"title": "RAFT \ud559\uc2b5 \ubc29\ubc95\ub860", "badge": "RAFT",
             "b1": "Oracle + Distractor \ubb38\uc11c \uad6c\ubd84 \ud559\uc2b5",
             "b2": "\ud658\uac01 \ucd5c\uc18c\ud654 + \uadfc\uac70 \uae30\ubc18 \ub2f5\ubcc0 \uc0dd\uc131",
             "b3": "\uc81c\ud488\ubcc4 QLoRA \uc5b4\ub311\ud130 \ub3d9\uc801 \ub85c\ub529"},
            {"title": "\uae30\uc5c5 \uc804\ud658 \ube44\uc804", "badge": "Vision",
             "b1": "Middleware Vendor \u2192 AI Platform Company",
             "b2": "AI \uc11c\ube44\uc2a4 \uad6c\ub3c5(SaaS) \ub9e4\ucd9c \ucc3d\ucd9c",
             "b3": "\uc608\uce21\ud615(Proactive) \uc2dc\uc2a4\ud15c \uad00\ub9ac \uc2e4\ud604"},
        ],
    )
    print(" 18  Why We Win (6 cards)")

    # 19. Future Roadmap (HTML slide 20)
    build_content(
        prs, "05", 19, "Future Roadmap",
        "\uc9c0\uc18d \uac00\ub2a5\ud55c",
        "\ud601\uc2e0\uc744 \uc704\ud55c \uccad\uc0ac\uc9c4",
        "AI \uc9c0\ub2a5\ud615 \uac80\uc0c9\uc5d0\uc11c \uc2dc\uc791\ud558\uc5ec "
        "\uc790\uc728 \ub9c8\uc774\uadf8\ub808\uc774\uc158, "
        "\uae00\ub85c\ubc8c AI SaaS \uc0dd\ud0dc\uacc4\ub85c \ud655\uc7a5\ud569\ub2c8\ub2e4.",
        "Phase 1 (Current): AI KMS",
        "AI \uc9c0\ub2a5\ud615 \uac80\uc0c9 \ubc0f \uae30\uc220 \uc9c0\uc6d0 \uc790\ub3d9\ud654. "
        "Agentic RAG + QLoRA \ub3c4\uba54\uc778 \ud2b9\ud654 \ud559\uc2b5 \uae30\ubc18 "
        "OpenFrame \uc804\uc6a9 \uc9c0\uc2dd\uad00\ub9ac\uc2dc\uc2a4\ud15c \uc6b4\uc601.",
        "Phase 2 (2026): Autonomous Migration Agent",
        "\uc790\uc728 \ub9c8\uc774\uadf8\ub808\uc774\uc158 \uc5d0\uc774\uc804\ud2b8 \uac1c\ubc1c. "
        "\ucf54\ub4dc \uc790\ub3d9 \ubcc0\ud658 \ubc0f \ud14c\uc2a4\ud2b8 \uc2dc\ub098\ub9ac\uc624 \uc790\ub3d9 \uc0dd\uc131. "
        "AS-IS \uc790\uc0b0 \ubd84\uc11d\ubd80\ud130 \uc804\ud658 \ub85c\ub4dc\ub9f5 \uc81c\uc548\uae4c\uc9c0 "
        "End-to-End AI \ub9c8\uc774\uadf8\ub808\uc774\uc158 \uc9c0\uc6d0.",
        "Phase 3 (2027+): Global AI SaaS",
        "OpenFrame AI API \uacf5\uac1c \ubc0f \ud30c\ud2b8\ub108 \ud50c\ub7ec\uadf8\uc778 \uc0dd\ud0dc\uacc4 \uad6c\ucd95",
        "\ud574\uc678 \uc2dc\uc7a5 \ub300\uc0c1 \uc644\uc804 \uad00\ub9ac\ud615 AI Modernization \uc11c\ube44\uc2a4",
    )
    print(" 19  Future Roadmap")

    # 20. Key Messaging (HTML slide 22)
    build_content(
        prs, "05", 20, "Key Messaging",
        "Tagline &",
        "Elevator Pitch",
        "\uacbd\uc601\uc9c4\uacfc \uace0\uac1d\uc5d0\uac8c \uc804\ub2ec\ud560 "
        "\ud575\uc2ec \uba54\uc2dc\uc9c0\ub97c \uc815\ub9ac\ud569\ub2c8\ub2e4.",
        "Tagline Candidates",
        "1. 'Legacy Intelligence, Modernized Instantly.' "
        "\u2014 \ub808\uac70\uc2dc \uc9c0\uc2dd\uc744 \uc989\uc2dc \ud604\ub300\ud654. "
        "2. 'Your AI Partner for Mainframe Modernization.' "
        "\u2014 \uba54\uc778\ud504\ub808\uc784 \ud604\ub300\ud654\uc758 AI \ud30c\ud2b8\ub108. "
        "3. 'Unlock the Blackbox: OpenFrame AI.' "
        "\u2014 \ube14\ub799\ubc15\uc2a4\ub97c \uc5f4\uc5b4\ub77c.",
        "Elevator Pitch (30\ucd08)",
        "\uae30\uc5c5\uc758 \ud575\uc2ec \ub808\uac70\uc2dc \uc2dc\uc2a4\ud15c\uc740 \uc804\ubb38 \uc778\ub825 \ubd80\uc871\uacfc "
        "\ube14\ub799\ubc15\uc2a4\ud654\ub85c \uc704\uae30\uc5d0 \ucc98\ud574 \uc788\uc2b5\ub2c8\ub2e4. OpenFrame AI KMS\ub294 "
        "\uae30\uc5c5\uc758 \uc218\uc2ed \ub144 \ub41c \uae30\uc220 \uc790\uc0b0\uc744 \uc644\ubcbd\ud788 \uc774\ud574\ud558\uace0 "
        "\uad6c\uc870\uc801\uc73c\ub85c \ubd84\uc11d\ud558\ub294 '\ub3c4\uba54\uc778 \ud2b9\ud654 AI \uc5d4\uc9c0\ub2c8\uc5b4\ub9c1 \ud50c\ub7ab\ud3fc'\uc785\ub2c8\ub2e4.",
        "Core Message",
        "\ub2e8\uc21c \ucc57\ubd07\uc774 \uc544\ub2cc \uc5d4\uc9c0\ub2c8\uc5b4\ub9c1 \ud50c\ub7ab\ud3fc\uc73c\ub85c \ucc28\ubcc4\ud654",
        "Re-inventing Legacy with AI \u2014 AI\ub85c \ub808\uac70\uc2dc\ub97c \uc7ac\ud0c4\uc0dd",
    )
    print(" 20  Key Messaging")

    # 21. Thank You (HTML slide 23)
    build_closing(prs)
    print(" 21  Thank You")

    # Delete original 5 template slides
    for _ in range(5):
        delete_slide(prs, 0)

    prs.save(OUTPUT_PATH)
    print(f"\nGenerated: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
