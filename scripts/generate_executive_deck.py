"""
TmaxSoft AI KMS Executive Marketing Deck Generator
Generates a professional PowerPoint presentation from the executive deck content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# DESIGN SYSTEM
# ============================================================
# Colors
BG_DARK = RGBColor(0x0A, 0x0E, 0x17)
BG_CARD = RGBColor(0x11, 0x18, 0x27)
BG_ACCENT = RGBColor(0x1A, 0x22, 0x34)
BORDER = RGBColor(0x2A, 0x35, 0x48)
TEXT_PRIMARY = RGBColor(0xE6, 0xED, 0xF3)
TEXT_SECONDARY = RGBColor(0xC9, 0xD1, 0xD9)
TEXT_DIM = RGBColor(0x6B, 0x7B, 0x8F)
BLUE = RGBColor(0x4A, 0x9E, 0xFF)
GREEN = RGBColor(0x34, 0xD3, 0x99)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
CYAN = RGBColor(0x67, 0xE8, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)

# Fonts
FONT_TITLE = 'Segoe UI'
FONT_BODY = 'Segoe UI'
FONT_MONO = 'Cascadia Code'

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Minimal corner radius
    shape.adjustments[0] = 0.02
    return shape


def add_accent_bar(slide, left, top, height, color=BLUE, width=Pt(4)):
    """Add a colored accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text='', font_size=14,
                 font_color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=None):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name or FONT_BODY
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    font_color=TEXT_SECONDARY, bullet_color=None, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_BODY
        p.space_after = spacing
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
    return txBox


def add_quote(slide, left, top, width, text, font_size=14):
    """Add a styled quote block."""
    # Quote bar
    add_accent_bar(slide, left, top + Pt(4), Pt(font_size * 2.5 + 10), BLUE)
    # Quote text
    txBox = add_text_box(
        slide, left + Pt(16), top, width - Pt(16), Inches(1),
        text, font_size=font_size, font_color=BLUE, bold=False
    )
    txBox.text_frame.paragraphs[0].font.italic = True
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None,
              header_color=BG_ACCENT, row_color=BG_CARD):
    """Add a styled table."""
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows > 0 else 1
    row_h = Pt(32)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)

            # Cell formatting
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = FONT_BODY
                if r_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_PRIMARY
                else:
                    p.font.color.rgb = TEXT_SECONDARY

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r_idx == 0:
                cell_fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
            else:
                cell_fill.fore_color.rgb = row_color

            # Cell margins
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)

    return table_shape


def add_badge(slide, left, top, text, bg_color, text_color, width=None):
    """Add a small badge/pill."""
    w = width or Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Pt(24))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_slide_number(slide, num, total):
    """Add slide number in bottom right."""
    add_text_box(
        slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.3),
        f'{num} / {total}', font_size=8, font_color=TEXT_DIM,
        alignment=PP_ALIGN.RIGHT
    )


def add_divider_line(slide, left, top, width):
    """Add a thin horizontal divider."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()


# ============================================================
# SLIDE BUILDERS
# ============================================================
TOTAL_SLIDES = 25


def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Top accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # Centered content
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(0.5),
                 'TmaxSoft Japan', font_size=16, font_color=TEXT_DIM, bold=False)

    add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(1.2),
                 'OpenFrame AI KMS Platform', font_size=44, font_color=WHITE, bold=True)

    add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
                 'Domain-Trained AI for Legacy Modernization Intelligence',
                 font_size=20, font_color=BLUE, bold=False)

    add_divider_line(slide, Inches(1.5), Inches(4.8), Inches(3))

    add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.5),
                 '"The AI That Speaks Mainframe."',
                 font_size=16, font_color=TEXT_DIM, bold=False)

    add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
                 'Executive Strategy Briefing  |  Confidential',
                 font_size=11, font_color=TEXT_DIM)

    add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
                 '2026', font_size=11, font_color=TEXT_DIM)


def slide_02_problem(prs):
    """The $47 Billion Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'EXECUTIVE PROBLEM STATEMENT', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'The $47 Billion Problem', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
                 'Legacy modernization is failing — because of knowledge, not technology.',
                 font_size=15, font_color=BLUE)

    # Stat boxes row
    stats = [
        ('$47B', 'Global legacy\nmodernization market', BLUE),
        ('73%', 'COBOL systems\nstill mission-critical', GREEN),
        ('58+', 'Average mainframe\nengineer age', ORANGE),
        ('68%', 'Projects delayed by\nundocumented knowledge', RED),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)
        card = add_shape(slide, x, Inches(2.4), Inches(2.8), Inches(1.6), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(2.4), Inches(1.6), color)
        add_text_box(slide, x + Pt(16), Inches(2.55), Inches(2.4), Inches(0.6),
                     num, font_size=36, font_color=color, bold=True)
        add_text_box(slide, x + Pt(16), Inches(3.2), Inches(2.4), Inches(0.7),
                     label, font_size=11, font_color=TEXT_DIM)

    # Bullet points
    bullets = [
        'Incident resolution depends on tribal knowledge that exists in no system',
        'Average resolution time for complex ABEND: 4\u20138 hours of senior engineer time',
        'New engineers require 18\u201324 months to become independently productive',
        'Cost of a single expert departure: \u00a515M\u201330M in lost productivity over 12 months',
    ]
    add_bullet_list(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5),
                    bullets, font_size=13, font_color=TEXT_SECONDARY)

    add_quote(slide, Inches(0.8), Inches(6.5), Inches(11),
              '"The technology to migrate exists. The knowledge to migrate correctly is disappearing."',
              font_size=13)

    add_slide_number(slide, 2, TOTAL_SLIDES)


def slide_03_knowledge_fragmentation(prs):
    """Knowledge Fragmentation: The Silent Risk."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'THE HIDDEN COST', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Knowledge Fragmentation: The Silent Risk', font_size=32, font_color=WHITE, bold=True)

    # Left column - problems
    add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5), BG_CARD, BORDER)
    add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(4.5), RED)
    add_text_box(slide, Inches(1.1), Inches(2.15), Inches(5.3), Inches(0.4),
                 'Current State', font_size=14, font_color=RED, bold=True)

    problems = [
        '245+ technical manuals across 19 products (PDF-only)',
        'Error resolution procedures exist only in engineer memory',
        'Cross-product dependency chains require holistic understanding',
        'Average incident resolution: 4\u20138 hours for complex ABEND',
        'New engineers need 18\u201324 months of training',
        '60\u201370% of legacy systems lack current documentation',
    ]
    add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.3), Inches(3.5),
                    problems, font_size=12, font_color=TEXT_SECONDARY, spacing=Pt(10))

    # Right column - impact
    add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), BG_CARD, BORDER)
    add_accent_bar(slide, Inches(7.0), Inches(2.0), Inches(4.5), ORANGE)
    add_text_box(slide, Inches(7.3), Inches(2.15), Inches(5.0), Inches(0.4),
                 'Business Impact', font_size=14, font_color=ORANGE, bold=True)

    impacts = [
        '\u00a515M\u201330M  lost productivity per expert departure',
        '\u00a52.4M+  annual cost per major migration project overrun',
        '3\u20136 months  delay when design errors found late',
        '60\u201370%  of L1 tickets escalate unnecessarily',
        'After-hours incidents wait until expert is available',
        'Knowledge walks out the door with every retirement',
    ]
    add_bullet_list(slide, Inches(7.3), Inches(2.7), Inches(5.0), Inches(3.5),
                    impacts, font_size=12, font_color=TEXT_SECONDARY, spacing=Pt(10))

    add_quote(slide, Inches(0.8), Inches(6.7), Inches(11),
              '"Every expert departure is a silent data breach \u2014 of knowledge."', font_size=13)

    add_slide_number(slide, 3, TOTAL_SLIDES)


def slide_04_why_generic_fails(prs):
    """Why Generic RAG Cannot Solve This."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'THE GAP', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Why Generic RAG Cannot Solve This', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
                 'Standard AI retrieval fails on domain-specific structural reasoning.',
                 font_size=15, font_color=RED)

    # Comparison table
    data = [
        ['Challenge', 'Generic RAG', 'Required Capability'],
        ['"What is ABEND S0C7?"', 'Guesses from general knowledge', 'Exact error registry lookup + context'],
        ['"Parse this JCL"', 'Treats as plain text', 'Structural JOB/EXEC/DD parsing'],
        ['"tjesmgr BOOT fails"', 'No knowledge of TJES', '19-product knowledge base retrieval'],
        ['"Compare OSC vs CICS"', 'Generic comparison', 'Product-specific architecture analysis'],
        ['"Diagnose batch failure"', 'Single-turn response', 'Multi-agent diagnosis pipeline'],
    ]
    add_table(slide, Inches(0.8), Inches(2.3), Inches(11.7), data,
              col_widths=[Inches(2.8), Inches(4.2), Inches(4.7)])

    bullets = [
        'Generic RAG treats all text equally \u2014 cannot parse JCL/COBOL structures',
        'OpenAI/ChatGPT has zero training data on OpenFrame, TJES, TACF, or AIM/DB',
        'No hallucination detection \u2014 plausible but wrong answers are worse than no answer',
        'Cannot perform multi-step diagnostic reasoning (parse \u2192 identify \u2192 trace \u2192 fix)',
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.5),
                    bullets, font_size=12)

    add_quote(slide, Inches(0.8), Inches(6.5), Inches(11),
              '"Hallucination in enterprise support is not an inconvenience. It is a liability."',
              font_size=13)

    add_slide_number(slide, 4, TOTAL_SLIDES)


def slide_05_why_now(prs):
    """Market Timing: Why Now."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'MARKET TIMING', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Why Now: The Strategic Window', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
                 'Three irreversible forces are converging \u2014 and TmaxSoft is positioned to own this moment.',
                 font_size=15, font_color=BLUE)

    forces = [
        ('Force 1', 'AI Transformation Mandate', 'Every enterprise board has AI on the agenda.\nLegacy IT must participate or be replaced.', BLUE),
        ('Force 2', 'Knowledge Cliff 2025\u20132030', 'The retirement wave is not a prediction.\nIt is demographic certainty.', ORANGE),
        ('Force 3', 'Global Competition', 'IBM, Micro Focus, and hyperscalers investing.\nDomain-specialist window: 3\u20135 years.', RED),
        ('Force 4', 'Customer Expectation Shift', 'Enterprise customers now expect AI-first\nsupport, not ticket-based workflows.', GREEN),
    ]

    for i, (label, title, desc, color) in enumerate(forces):
        x = Inches(0.8 + i * 3.1)
        add_shape(slide, x, Inches(2.4), Inches(2.8), Inches(2.8), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(2.4), Inches(2.8), color)
        add_text_box(slide, x + Pt(16), Inches(2.55), Inches(2.4), Inches(0.3),
                     label, font_size=10, font_color=color, bold=True)
        add_text_box(slide, x + Pt(16), Inches(2.9), Inches(2.4), Inches(0.5),
                     title, font_size=16, font_color=WHITE, bold=True)
        add_text_box(slide, x + Pt(16), Inches(3.5), Inches(2.4), Inches(1.2),
                     desc, font_size=11, font_color=TEXT_DIM)

    add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.4), BG_ACCENT, BORDER)
    add_text_box(slide, Inches(1.2), Inches(5.75), Inches(10.8), Inches(0.4),
                 'Strategic Implication', font_size=13, font_color=BLUE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(6.15), Inches(10.8), Inches(0.6),
                 'TmaxSoft can either be the AI platform that powers modernization\nor a product that competes on features. This is a positioning decision, not a product decision.',
                 font_size=13, font_color=TEXT_PRIMARY)

    add_slide_number(slide, 5, TOTAL_SLIDES)


def slide_06_product_overview(prs):
    """Introducing OpenFrame AI KMS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'PRODUCT OVERVIEW', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Introducing OpenFrame AI KMS', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
                 'A domain-trained AI engineering platform that transforms fragmented knowledge into actionable intelligence.',
                 font_size=14, font_color=BLUE)

    # What It Is - 5 capability cards
    capabilities = [
        ('RAFT-Trained LLM', 'Qwen 32B with 24 product-\nspecific QLoRA adapters', PURPLE),
        ('Hybrid RAG Engine', 'Graph + Vector + BM25 fusion\nwith cross-encoder reranking', GREEN),
        ('Multi-Agent System', 'DAG-based parallel agents\nfor complex analysis', PINK),
        ('Embedded Parsers', 'JCL, COBOL, ASM structural\nanalyzers (deterministic)', CYAN),
        ('Anti-Hallucination', 'Sentence-level verification\nagainst source docs (95%+)', ORANGE),
    ]

    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(3), Inches(0.3),
                 'WHAT IT IS', font_size=11, font_color=BLUE, bold=True)

    for i, (title, desc, color) in enumerate(capabilities):
        x = Inches(0.8 + i * 2.45)
        add_shape(slide, x, Inches(2.6), Inches(2.2), Inches(1.8), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(2.6), Inches(1.8), color)
        add_text_box(slide, x + Pt(14), Inches(2.75), Inches(1.9), Inches(0.4),
                     title, font_size=12, font_color=color, bold=True)
        add_text_box(slide, x + Pt(14), Inches(3.2), Inches(1.9), Inches(0.9),
                     desc, font_size=10, font_color=TEXT_DIM)

    # What It Is NOT
    add_text_box(slide, Inches(0.8), Inches(4.7), Inches(3), Inches(0.3),
                 'WHAT IT IS NOT', font_size=11, font_color=RED, bold=True)

    nots = [
        ('Not a Chatbot', 'Structural AI engineering\nplatform with parsers'),
        ('Not Document Search', 'Knowledge graph with\n42,596 linked chunks'),
        ('Not a Generic AI Wrapper', 'Domain-trained with\nRAFT methodology'),
    ]
    for i, (title, desc) in enumerate(nots):
        x = Inches(0.8 + i * 4.0)
        add_shape(slide, x, Inches(5.1), Inches(3.7), Inches(1.2), RGBColor(0x1A, 0x0E, 0x0E), RGBColor(0x3D, 0x20, 0x20))
        add_text_box(slide, x + Inches(0.2), Inches(5.2), Inches(3.3), Inches(0.35),
                     '\u2717  ' + title, font_size=13, font_color=RED, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(5.55), Inches(3.3), Inches(0.6),
                     desc, font_size=11, font_color=TEXT_DIM)

    add_quote(slide, Inches(0.8), Inches(6.55), Inches(11),
              '"The first AI system that understands OpenFrame at a structural level."', font_size=13)

    add_slide_number(slide, 6, TOTAL_SLIDES)


def slide_07_architecture(prs):
    """Technology Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'TECHNOLOGY', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Purpose-Built AI Architecture', font_size=32, font_color=WHITE, bold=True)

    layers = [
        ('Layer 3 \u2014 Automation Engine', PINK, [
            'DAG-based multi-agent orchestration',
            '5-agent JCL diagnosis pipeline',
            'Vision LLM for PDF image/table analysis',
            'Faithfulness verification (cosine similarity)',
        ]),
        ('Layer 2 \u2014 Domain-Trained AI', PURPLE, [
            'Qwen 32B + RAFT methodology (CPT \u2192 SFT \u2192 DPO)',
            '24 QLoRA adapters on vLLM (continuous batching)',
            '32K context window \u00b7 A100 48GB \u00d7 4',
            'Product routing with 95%+ accuracy',
        ]),
        ('Layer 1 \u2014 Knowledge Foundation', GREEN, [
            '42,596 document chunks in Neo4j (Graph + Vector)',
            '13,450 domain entities (commands, configs, errors, APIs)',
            '245 technical manuals across 19 products',
            'Two-stage retrieval: Summary (<10ms) \u2192 Deep search',
        ]),
    ]

    for i, (title, color, items) in enumerate(layers):
        y = Inches(1.8 + i * 1.85)
        add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.65), BG_CARD, BORDER)
        add_accent_bar(slide, Inches(0.8), y, Inches(1.65), color)
        add_text_box(slide, Inches(1.2), y + Pt(8), Inches(10), Inches(0.35),
                     title, font_size=14, font_color=color, bold=True)
        for j, item in enumerate(items):
            add_text_box(slide, Inches(1.2 + (j % 2) * 5.5), y + Inches(0.45 + (j // 2) * 0.38),
                         Inches(5.3), Inches(0.35),
                         '\u2022  ' + item, font_size=11, font_color=TEXT_SECONDARY)

    add_slide_number(slide, 7, TOTAL_SLIDES)


def slide_08_comparison(prs):
    """Generic RAG vs AI KMS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'COMPETITIVE ANALYSIS', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Generic RAG vs TmaxSoft AI KMS', font_size=32, font_color=WHITE, bold=True)

    data = [
        ['Dimension', 'Generic RAG', 'TmaxSoft AI KMS'],
        ['Knowledge Base', 'General web data', '245 manuals, 19 products, 42K chunks'],
        ['LLM Training', 'Generic pre-training', 'RAFT + 24 QLoRA domain adapters'],
        ['Document Understanding', 'Text embedding only', 'Graph + Vector + BM25 hybrid'],
        ['Structural Parsing', 'None', 'JCL / COBOL / ASM parsers'],
        ['Error Diagnosis', 'Pattern matching', 'Multi-agent pipeline + ABEND registry'],
        ['Hallucination Control', 'None / basic', 'Sentence-level verification (95%+)'],
        ['Multi-step Reasoning', 'Single query-response', 'DAG orchestration, parallel agents'],
        ['Product Coverage', 'Generic', '19 OpenFrame products, version-specific'],
        ['Vision Analysis', 'None', 'PDF image/chart/table extraction'],
        ['Support Integration', 'Separate system', 'AI-first + human escalation built-in'],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), data,
              col_widths=[Inches(2.5), Inches(4.2), Inches(5.0)])

    add_quote(slide, Inches(0.8), Inches(6.7), Inches(11),
              '"The gap is not incremental. It is architectural."', font_size=14)

    add_slide_number(slide, 8, TOTAL_SLIDES)


def slide_09_business_impact(prs):
    """Quantified Business Impact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'BUSINESS IMPACT', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Quantified Business Impact', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
                 'Measurable, defensible improvements across every dimension of legacy operations.',
                 font_size=14, font_color=GREEN)

    data = [
        ['Metric', 'Current State', 'With AI KMS', 'Improvement'],
        ['Technical query resolution', '4\u20138 hours', '30\u201360 minutes', '70\u201390% faster'],
        ['Log / dump analysis', '2\u20134 hours (expert)', '10\u201330 minutes (AI)', '60\u201380% faster'],
        ['Migration design analysis', '2\u20133 weeks / module', '3\u20135 days / module', '15\u201325% reduction'],
        ['New engineer onboarding', '18\u201324 months', '6\u20139 months', '50\u201360% faster'],
        ['Incident escalation rate', '60\u201370% to L2/L3', '25\u201335% to L2/L3', '50% reduction'],
        ['Documentation coverage', '30\u201340% of systems', '85\u201395% of systems', '2.5\u00d7 increase'],
    ]
    add_table(slide, Inches(0.8), Inches(2.2), Inches(11.7), data,
              col_widths=[Inches(3.0), Inches(2.8), Inches(2.8), Inches(3.1)])

    # Bottom stat
    add_shape(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2), BG_ACCENT, BLUE)
    add_text_box(slide, Inches(1.5), Inches(5.85), Inches(10), Inches(0.45),
                 'Annual Operational Savings', font_size=14, font_color=TEXT_DIM, bold=False)
    add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
                 '\u00a545M \u2013 \u00a5120M per enterprise', font_size=28, font_color=GREEN, bold=True)

    add_slide_number(slide, 9, TOTAL_SLIDES)


def slide_10_roi(prs):
    """ROI Scenario."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'ROI ANALYSIS', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'ROI Scenario: Enterprise Financial Institution', font_size=28, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                 'Major bank \u00b7 500+ COBOL batch programs \u00b7 3 OpenFrame environments \u00b7 12-person team',
                 font_size=12, font_color=TEXT_DIM)

    data = [
        ['Cost Category', 'Without AI KMS', 'With AI KMS', 'Savings'],
        ['Senior engineer hours (incident)', '\u00a572M/year', '\u00a522M/year', '\u00a550M'],
        ['Migration consulting', '\u00a5180M/project', '\u00a5135M/project', '\u00a545M'],
        ['Training & onboarding', '\u00a524M/year', '\u00a510M/year', '\u00a514M'],
        ['Incident downtime cost', '\u00a536M/year', '\u00a512M/year', '\u00a524M'],
        ['Total annual cost', '\u00a5312M', '\u00a5179M', '\u00a5133M'],
    ]
    add_table(slide, Inches(0.8), Inches(2.1), Inches(11.7), data,
              col_widths=[Inches(3.5), Inches(2.7), Inches(2.7), Inches(2.8)])

    # ROI result cards
    results = [
        ('Platform Investment', '\u00a536M/year', TEXT_DIM, BORDER),
        ('Net Annual Savings', '\u00a597M', GREEN, GREEN),
        ('Year 1 ROI', '269%', BLUE, BLUE),
    ]
    for i, (label, value, vcolor, bcolor) in enumerate(results):
        x = Inches(0.8 + i * 4.0)
        add_shape(slide, x, Inches(5.1), Inches(3.7), Inches(1.1), BG_CARD, bcolor)
        add_text_box(slide, x + Inches(0.2), Inches(5.2), Inches(3.3), Inches(0.35),
                     label, font_size=12, font_color=TEXT_DIM)
        add_text_box(slide, x + Inches(0.2), Inches(5.5), Inches(3.3), Inches(0.5),
                     value, font_size=28, font_color=vcolor, bold=True)

    add_quote(slide, Inches(0.8), Inches(6.5), Inches(11),
              '"This is not a cost center. It is an operational leverage platform."', font_size=13)

    add_slide_number(slide, 10, TOTAL_SLIDES)


def _capability_slide(prs, slide_num, label, title, subtitle, color,
                      problems, automations, impacts, quote_text):
    """Generic capability block slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 label, font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 title, font_size=28, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
                 subtitle, font_size=13, font_color=color)

    # Three columns
    cols = [
        ('Problem', RED, problems),
        ('AI Automation', color, automations),
        ('Business Impact', GREEN, impacts),
    ]
    for i, (col_title, col_color, items) in enumerate(cols):
        x = Inches(0.8 + i * 4.1)
        add_shape(slide, x, Inches(2.1), Inches(3.8), Inches(3.9), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(2.1), Inches(3.9), col_color)
        add_text_box(slide, x + Pt(14), Inches(2.2), Inches(3.4), Inches(0.35),
                     col_title, font_size=13, font_color=col_color, bold=True)
        add_bullet_list(slide, x + Pt(14), Inches(2.65), Inches(3.4), Inches(3.0),
                        items, font_size=11, font_color=TEXT_SECONDARY, spacing=Pt(8))

    add_quote(slide, Inches(0.8), Inches(6.3), Inches(11), quote_text, font_size=12)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)


def slide_11_capability_support(prs):
    _capability_slide(prs, 11,
        'CAPABILITY 1 OF 4', 'AI Technical Support',
        'Transform reactive ticket-based support into proactive AI-first intelligence.', BLUE,
        [
            '4\u20138 hours per complex query',
            '60\u201370% tickets escalate to L2/L3',
            'Knowledge scattered across 245 PDFs',
            'Tribal knowledge dependency',
        ],
        [
            '19 products \u00d7 24 adapters routing',
            'Summary pre-retrieval (<10ms)',
            'Verified responses with source citation',
            'Real-time SSE streaming with trace',
        ],
        [
            '70\u201390% faster resolution',
            '50% reduction in escalations',
            '24/7 availability without staffing',
            'Consistent quality for all engineers',
        ],
        '"Support transforms from cost center to competitive advantage."')


def slide_12_capability_log(prs):
    _capability_slide(prs, 12,
        'CAPABILITY 2 OF 4', 'AI Log & Core Analysis',
        'Automated JCL job failure diagnosis replaces hours of expert analysis.', ORANGE,
        [
            'JCL failures need JOB/EXEC/DD analysis',
            'Cross-referencing multiple manuals',
            'Single complex failure: 2\u20134 hours',
            'After-hours: delayed until expert available',
        ],
        [
            '5-Agent diagnosis pipeline',
            'Structural JCL parser (deterministic)',
            'ABEND code registry with resolution',
            'Automated HTML diagnosis report',
        ],
        [
            '60\u201380% faster log analysis',
            'Automated after-hours first-response',
            'Consistent diagnosis quality',
            'Knowledge captured for future use',
        ],
        '"From hours of expert guessing to minutes of AI-verified diagnosis."')


def slide_13_capability_migration(prs):
    _capability_slide(prs, 13,
        'CAPABILITY 3 OF 4', 'AI Migration Design',
        'AI-assisted analysis reduces project risk and accelerates design decisions.', PURPLE,
        [
            'Thousands of COBOL programs to analyze',
            'Vendor-specific extensions (IBM/Fujitsu)',
            'Late-found errors: 3\u20136 month delay',
            'Scarce Fujitsu/IBM-experienced engineers',
        ],
        [
            'Legacy code analysis (COBOL/JCL/ASM)',
            'Cross-vendor: IBM MVS, Fujitsu XSP/MSP',
            'Automated compatibility assessment',
            'AI-generated conversion strategy',
        ],
        [
            '15\u201325% timeline reduction',
            'Early detection of incompatibilities',
            'Reduced expert dependency',
            'Risk-assessed recommendations',
        ],
        '"Design errors caught by AI before coding — not after deployment."')


def slide_14_capability_docs(prs):
    _capability_slide(prs, 14,
        'CAPABILITY 4 OF 4', 'AI Asset Documentation',
        'Transform undocumented legacy systems into structured, searchable knowledge.',TEAL,
        [
            '60\u201370% systems lack documentation',
            'Existing docs outdated, PDF-only',
            'New hires need 18+ months training',
            'Manual documentation is expensive',
        ],
        [
            'PDF \u2192 structured summaries auto-gen',
            '13,450 entities in knowledge graph',
            'Multi-language: JP, KR, EN native',
            'Living documentation, auto-updating',
        ],
        [
            'Coverage: 30% \u2192 90%+',
            'Onboarding: 18mo \u2192 6\u20139mo',
            'Self-service knowledge access',
            'Organizational knowledge preserved',
        ],
        '"Documentation that writes itself \u2014 and never goes stale."')


def slide_15_premium_support(prs):
    """Premium Technical Support."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'PREMIUM SUPPORT', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'AI + Human Hybrid Support Model', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
                 'AI handles 70%+ of queries. Human experts focus on high-value cases.',
                 font_size=14, font_color=GREEN)

    data = [
        ['Tier', 'Handler', 'Coverage', 'Response Time'],
        ['Tier 0', 'AI KMS Automated', 'Standard queries, error lookups, commands', 'Instant (<30 sec)'],
        ['Tier 1', 'AI-Assisted Engineer', 'Complex queries + human review', '30 min \u2013 2 hours'],
        ['Tier 2', 'Senior Expert + AI', 'Critical incidents, architecture decisions', '2 \u2013 8 hours'],
        ['Tier 3', 'R&D Escalation', 'Product bugs, deep investigation', '1 \u2013 5 business days'],
    ]
    add_table(slide, Inches(0.8), Inches(2.2), Inches(11.7), data,
              col_widths=[Inches(1.5), Inches(2.8), Inches(4.5), Inches(2.9)])

    features = [
        ('AI-First Triage', 'Every query analyzed by AI\nbefore human involvement', BLUE),
        ('Screen Sharing', 'Real-time remote collaboration\non complex issues', GREEN),
        ('Auto Report Gen', 'AI produces structured\ndiagnosis for expert review', ORANGE),
        ('Knowledge Capture', 'Every resolution feeds back\ninto knowledge graph', PURPLE),
    ]
    for i, (title, desc, color) in enumerate(features):
        x = Inches(0.8 + i * 3.1)
        add_shape(slide, x, Inches(4.6), Inches(2.8), Inches(1.4), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(4.6), Inches(1.4), color)
        add_text_box(slide, x + Pt(14), Inches(4.7), Inches(2.4), Inches(0.35),
                     title, font_size=12, font_color=color, bold=True)
        add_text_box(slide, x + Pt(14), Inches(5.1), Inches(2.4), Inches(0.7),
                     desc, font_size=11, font_color=TEXT_DIM)

    add_quote(slide, Inches(0.8), Inches(6.3), Inches(11),
              '"Support is not a cost center. With AI, it becomes a competitive moat."', font_size=13)

    add_slide_number(slide, 15, TOTAL_SLIDES)


def slide_16_strategic_position(prs):
    """The AI Operating Layer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'STRATEGIC POSITIONING', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'The AI Operating Layer for OpenFrame', font_size=32, font_color=WHITE, bold=True)

    # Current state
    add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.0), BG_CARD, BORDER)
    add_text_box(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.35),
                 'Current State (Product-Centric)', font_size=13, font_color=RED, bold=True)
    add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5.0), Inches(1.3),
                 'Customer \u2192 OpenFrame Product \u2192 Manual Support \u2192 Manual Docs\n\n'
                 '\u2022 Linear support model\n'
                 '\u2022 Knowledge trapped in silos\n'
                 '\u2022 Each interaction is isolated',
                 font_size=11, font_color=TEXT_SECONDARY, font_name=FONT_MONO)

    # Future state
    add_shape(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(2.0), BG_CARD, BORDER)
    add_accent_bar(slide, Inches(6.8), Inches(2.0), Inches(2.0), BLUE)
    add_text_box(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.35),
                 'Future State (AI Platform)', font_size=13, font_color=BLUE, bold=True)
    add_text_box(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(1.3),
                 'Customer \u2192 AI KMS Layer \u2192 Products + AI + Auto-Support\n\n'
                 '\u2022 Every touchpoint is intelligent\n'
                 '\u2022 Knowledge compounds over time\n'
                 '\u2022 Platform gets smarter with usage',
                 font_size=11, font_color=TEXT_SECONDARY, font_name=FONT_MONO)

    # Arrow
    add_text_box(slide, Inches(5.9), Inches(2.7), Inches(1.2), Inches(0.5),
                 '\u25B6', font_size=32, font_color=BLUE, alignment=PP_ALIGN.CENTER)

    # Platform values
    values = [
        ('Product Intelligence', 'Every product enhanced\nwith AI understanding', PURPLE),
        ('Operational Intelligence', 'Support, diagnosis, migration\nall AI-augmented', GREEN),
        ('Knowledge Intelligence', 'Living graph grows\nwith every interaction', ORANGE),
        ('Ecosystem Intelligence', 'API layer for partner\nand customer integrations', CYAN),
    ]
    for i, (title, desc, color) in enumerate(values):
        x = Inches(0.8 + i * 3.1)
        add_shape(slide, x, Inches(4.5), Inches(2.8), Inches(1.4), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(4.5), Inches(1.4), color)
        add_text_box(slide, x + Pt(14), Inches(4.6), Inches(2.4), Inches(0.35),
                     title, font_size=12, font_color=color, bold=True)
        add_text_box(slide, x + Pt(14), Inches(5.0), Inches(2.4), Inches(0.7),
                     desc, font_size=11, font_color=TEXT_DIM)

    add_quote(slide, Inches(0.8), Inches(6.2), Inches(11),
              '"The AI layer transforms every touchpoint from manual to intelligent."', font_size=13)

    add_slide_number(slide, 16, TOTAL_SLIDES)


def slide_17_business_model(prs):
    """Business Model Transformation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'BUSINESS TRANSFORMATION', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'From Product Vendor to AI Platform Company', font_size=28, font_color=WHITE, bold=True)

    data = [
        ['Dimension', 'Current Model', 'AI Platform Model'],
        ['Revenue', 'License + maintenance', 'SaaS subscription + AI usage + premium'],
        ['Engagement', 'Project-based', 'Continuous platform relationship'],
        ['Value Delivery', 'One-time migration', 'Ongoing AI intelligence'],
        ['Differentiation', 'Product features', 'AI capability + domain knowledge'],
        ['Scalability', 'Linear (per project)', 'Exponential (per user, per query)'],
        ['Growth Vector', 'New customers only', 'Expansion within accounts + new verticals'],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), data,
              col_widths=[Inches(2.5), Inches(4.2), Inches(5.0)])

    # Revenue projection
    add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), BG_ACCENT, BORDER)
    add_text_box(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.35),
                 'Revenue Impact Projection', font_size=13, font_color=BLUE, bold=True)

    projections = [
        ('Year 1', '\u00a5300M', '5\u20138 enterprise accounts (early adopters)', BLUE),
        ('Year 2', '\u00a5750M', 'Expansion + new segments', GREEN),
        ('Year 3', '\u00a51.5B', 'Platform maturity + global markets', PURPLE),
    ]
    for i, (year, amount, desc, color) in enumerate(projections):
        x = Inches(1.2 + i * 3.8)
        add_text_box(slide, x, Inches(5.7), Inches(1.5), Inches(0.3),
                     year, font_size=11, font_color=TEXT_DIM, bold=True)
        add_text_box(slide, x, Inches(6.0), Inches(2), Inches(0.45),
                     amount, font_size=24, font_color=color, bold=True)
        add_text_box(slide, x, Inches(6.4), Inches(3.5), Inches(0.35),
                     desc, font_size=10, font_color=TEXT_DIM)

    add_slide_number(slide, 17, TOTAL_SLIDES)


def slide_18_competitive_moat(prs):
    """Competitive Moat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'COMPETITIVE MOAT', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Four Reinforcing Barriers', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
                 'Defensible advantages that widen over time.', font_size=14, font_color=BLUE)

    moats = [
        ('Moat 1', 'Domain Training Data', '245 manuals, 42,596 chunks, 13,450 entities.\nCompetitors need years to replicate.', PURPLE,
         'Barrier: Data takes time. Knowledge takes experience.'),
        ('Moat 2', 'Structural Parsers', 'JCL, COBOL, ASM parsers built into platform.\nDeterministic analysis, not AI inference.', GREEN,
         'Barrier: Requires deep product knowledge.'),
        ('Moat 3', 'Product Adapters', '24 QLoRA adapters, each product-trained.\n95%+ routing accuracy, version-specific.', ORANGE,
         'Barrier: Improves with each deployment. Network effect.'),
        ('Moat 4', 'Feedback Loop', 'Every query feeds back into training.\nKnowledge graph grows continuously.', BLUE,
         'Barrier: First-mover knowledge accumulation.'),
    ]

    for i, (label, title, desc, color, barrier) in enumerate(moats):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.2 + (i // 2) * 2.5)
        add_shape(slide, x, y, Inches(5.8), Inches(2.2), BG_CARD, BORDER)
        add_accent_bar(slide, x, y, Inches(2.2), color)
        add_text_box(slide, x + Pt(14), y + Pt(8), Inches(1.5), Inches(0.25),
                     label, font_size=10, font_color=color, bold=True)
        add_text_box(slide, x + Pt(14), y + Inches(0.35), Inches(5.3), Inches(0.35),
                     title, font_size=16, font_color=WHITE, bold=True)
        add_text_box(slide, x + Pt(14), y + Inches(0.7), Inches(5.3), Inches(0.7),
                     desc, font_size=11, font_color=TEXT_SECONDARY)
        add_text_box(slide, x + Pt(14), y + Inches(1.5), Inches(5.3), Inches(0.5),
                     barrier, font_size=10, font_color=TEXT_DIM, bold=False)

    add_slide_number(slide, 18, TOTAL_SLIDES)


def slide_19_why_not_generic(prs):
    """Why Not Generic AI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'DIFFERENTIATION', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Why Not Generic AI?', font_size=32, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
                 'Test: "Diagnose JCL job failure with ABEND S0C7 in STEP03"',
                 font_size=13, font_color=TEXT_DIM, font_name=FONT_MONO)

    data = [
        ['Approach', 'What Happens', 'Result'],
        ['ChatGPT / GPT-4', 'Generic mainframe advice.\nNo OpenFrame context.', '\u274c Plausible but useless'],
        ['Generic RAG', 'Text chunks mentioning S0C7.\nNo structural understanding.', '\u274c Partial, no diagnosis'],
        ['AI KMS', 'Parse JCL \u2192 Identify STEP03 \u2192\nLookup ABEND \u2192 Check DD \u2192 Report', '\u2705 Complete, verified'],
    ]
    add_table(slide, Inches(0.8), Inches(2.2), Inches(11.7), data,
              col_widths=[Inches(2.5), Inches(5.5), Inches(3.7)])

    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(10), Inches(0.35),
                 'Why the gap exists:', font_size=14, font_color=ORANGE, bold=True)

    reasons = [
        'Generic AI has zero OpenFrame training data',
        'Generic RAG has no structural parsing capability',
        'Neither can perform multi-step diagnostic reasoning',
        'Neither has domain-specific hallucination detection',
    ]
    add_bullet_list(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(2.0),
                    reasons, font_size=13, font_color=TEXT_SECONDARY)

    add_quote(slide, Inches(0.8), Inches(6.4), Inches(11),
              '"You cannot prompt-engineer domain expertise. You must train it."', font_size=14)

    add_slide_number(slide, 19, TOTAL_SLIDES)


def slide_20_raft(prs):
    """RAFT Training Methodology."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'AI SCIENCE', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'RAFT: Domain Training Methodology', font_size=28, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
                 'Cornell University Research (arXiv:2403.10131) \u2014 adapted for OpenFrame domain.',
                 font_size=12, font_color=TEXT_DIM)

    # Three phases
    phases = [
        ('Phase 1: CPT', 'Continued Pre-Training', 'Domain knowledge injection\n72MB raw text, 34.3M tokens\nQwen 32B base model\nFSDP GPU 4\u20137, 2 epochs', PURPLE,
         'Eval Perplexity: 1.65'),
        ('Phase 2: SFT', 'Supervised Fine-Tuning', '22 product-specific adapters\nChatML format, LoRA r=64\nQwen 7B \u00d7 22 products\n4 adapters parallel training', GREEN,
         '69 minutes total'),
        ('Phase 3: DPO', 'Preference Optimization', '2,000 preference pairs\nChosen vs rejected responses\nCross-product distractors\nRAFT oracle/distractor pattern', ORANGE,
         'Accuracy: 95%'),
    ]
    for i, (label, title, desc, color, metric) in enumerate(phases):
        x = Inches(0.8 + i * 4.1)
        add_shape(slide, x, Inches(2.1), Inches(3.8), Inches(3.2), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(2.1), Inches(3.2), color)
        add_text_box(slide, x + Pt(14), Inches(2.2), Inches(3.4), Inches(0.3),
                     label, font_size=11, font_color=color, bold=True)
        add_text_box(slide, x + Pt(14), Inches(2.5), Inches(3.4), Inches(0.4),
                     title, font_size=16, font_color=WHITE, bold=True)
        add_text_box(slide, x + Pt(14), Inches(3.0), Inches(3.4), Inches(1.5),
                     desc, font_size=11, font_color=TEXT_SECONDARY)
        # Metric badge
        add_shape(slide, x + Inches(0.2), Inches(4.6), Inches(3.2), Inches(0.45),
                  RGBColor(0x0D, 0x11, 0x17), color)
        add_text_box(slide, x + Inches(0.2), Inches(4.62), Inches(3.2), Inches(0.4),
                     metric, font_size=12, font_color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Anti-hallucination result
    add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), BG_ACCENT, BORDER)
    add_text_box(slide, Inches(1.2), Inches(5.7), Inches(10), Inches(0.35),
                 'Anti-Hallucination Verification', font_size=13, font_color=BLUE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10), Inches(0.5),
                 'E2E: 45 test cases across 8 components  \u00b7  Sentence-level cosine similarity  \u00b7  Faithfulness: 95%+ verified  \u00b7  <5% hallucination rate',
                 font_size=12, font_color=TEXT_SECONDARY)

    add_slide_number(slide, 20, TOTAL_SLIDES)


def slide_21_future_vision(prs):
    """Future Vision: Autonomous Migration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'FUTURE VISION', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Roadmap: Toward Autonomous Migration', font_size=32, font_color=WHITE, bold=True)

    phases = [
        ('Phase 1', 'Current', 'AI-Assisted\nIntelligence',
         'Knowledge retrieval +\ndiagnosis automation\nHuman-in-the-loop\n70\u201390% productivity gain', BLUE),
        ('Phase 2', '2026\u20132027', 'AI-Driven\nDesign',
         'Migration impact analysis\nAI conversion specs\nPredictive compatibility\nHuman review of AI designs', GREEN),
        ('Phase 3', '2027\u20132028', 'Autonomous\nMigration Agent',
         'End-to-end code conversion\nSelf-healing pipeline\nContinuous learning\nHuman oversight only', PURPLE),
        ('Phase 4', '2028+', 'Global AI\nPlatform',
         'Multi-tenant SaaS\nAPI ecosystem\nCross-platform support\nVertical solutions', ORANGE),
    ]

    for i, (label, timeline, title, desc, color) in enumerate(phases):
        x = Inches(0.8 + i * 3.1)
        add_shape(slide, x, Inches(2.0), Inches(2.8), Inches(4.5), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(2.0), Inches(4.5), color)
        add_text_box(slide, x + Pt(14), Inches(2.1), Inches(2.4), Inches(0.25),
                     label + '  \u00b7  ' + timeline, font_size=10, font_color=color, bold=True)
        add_text_box(slide, x + Pt(14), Inches(2.45), Inches(2.4), Inches(0.7),
                     title, font_size=18, font_color=WHITE, bold=True)
        add_text_box(slide, x + Pt(14), Inches(3.3), Inches(2.4), Inches(2.5),
                     desc, font_size=11, font_color=TEXT_SECONDARY)

        # Arrow between phases
        if i < 3:
            add_text_box(slide, Inches(3.4 + i * 3.1), Inches(3.8), Inches(0.5), Inches(0.5),
                         '\u25B6', font_size=18, font_color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # Current position indicator
    add_shape(slide, Inches(0.8), Inches(6.7), Inches(2.8), Inches(0.4),
              BLUE, BLUE)
    add_text_box(slide, Inches(0.8), Inches(6.72), Inches(2.8), Inches(0.35),
                 'WE ARE HERE', font_size=10, font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 21, TOTAL_SLIDES)


def slide_22_global_expansion(prs):
    """Global Expansion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'GLOBAL STRATEGY', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Global Market Expansion', font_size=32, font_color=WHITE, bold=True)

    data = [
        ['Market', 'Opportunity', 'Timeline', 'Strategy'],
        ['Japan', '\u00a52.1T installed base', '2025\u20132026', 'Direct enterprise. Lead market.'],
        ['Korea', '\u00a5800B modernization', '2026\u20132027', 'TmaxSoft Korea synergy.'],
        ['ASEAN', '\u00a5500B emerging IT', '2027\u20132028', 'Partner-led. Localized adapters.'],
        ['North America', '\u00a512T enterprise IT', '2027\u20132029', 'Cloud marketplace. SaaS.'],
        ['Europe', '\u00a54T regulated industry', '2028\u20132030', 'Compliance-focused. Banking.'],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), data,
              col_widths=[Inches(2.0), Inches(3.0), Inches(2.5), Inches(4.2)])

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.35),
                 'Scaling Advantage', font_size=14, font_color=BLUE, bold=True)

    advantages = [
        'Domain adapters are language-agnostic (product knowledge, not natural language)',
        'QLoRA adapters can be deployed per-region with local fine-tuning',
        'API ecosystem enables partner-built solutions on AI KMS platform',
        'Each market deployment compounds the knowledge graph',
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(1.8),
                    advantages, font_size=12)

    add_slide_number(slide, 22, TOTAL_SLIDES)


def slide_23_investment(prs):
    """Investment & Resources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'INVESTMENT', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Investment & Resource Plan', font_size=32, font_color=WHITE, bold=True)

    # Infrastructure and Team columns
    sections = [
        ('Infrastructure', BLUE, [
            'NVIDIA A100/H100 GPU cluster',
            'Neo4j Enterprise (Graph + Vector)',
            'vLLM continuous batching inference',
            'Kubernetes-ready containers',
        ]),
        ('Team', GREEN, [
            'AI/ML Engineers: LLM + adapters',
            'Domain Engineers: parsers + knowledge',
            'Platform Engineers: infra + scaling',
            'AI Product: feedback \u2192 training loop',
        ]),
    ]
    for i, (title, color, items) in enumerate(sections):
        x = Inches(0.8 + i * 6.2)
        add_shape(slide, x, Inches(1.8), Inches(5.8), Inches(2.5), BG_CARD, BORDER)
        add_accent_bar(slide, x, Inches(1.8), Inches(2.5), color)
        add_text_box(slide, x + Pt(14), Inches(1.9), Inches(5.3), Inches(0.35),
                     title, font_size=14, font_color=color, bold=True)
        add_bullet_list(slide, x + Pt(14), Inches(2.35), Inches(5.3), Inches(1.8),
                        items, font_size=12, font_color=TEXT_SECONDARY, spacing=Pt(8))

    # Investment timeline
    data = [
        ['', 'Year 1', 'Year 2', 'Year 3'],
        ['Investment', '\u00a5150M', '\u00a5250M', '\u00a5180M'],
        ['Focus', 'Infra + Core + Initial', 'Scale + Adapters + Markets', 'Optimize + SaaS + Global'],
        ['Projected Revenue', '\u00a5300M', '\u00a5750M', '\u00a51.5B'],
        ['Net Position', '\u00a5150M revenue', '\u00a5500M revenue', '\u00a51.32B revenue'],
    ]
    add_table(slide, Inches(0.8), Inches(4.6), Inches(11.7), data,
              col_widths=[Inches(2.8), Inches(3.0), Inches(3.0), Inches(2.9)])

    add_slide_number(slide, 23, TOTAL_SLIDES)


def slide_24_closing(prs):
    """Closing Executive Message."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'THE DECISIVE MOMENT', font_size=10, font_color=TEXT_DIM, bold=True)

    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
                 'Three Facts', font_size=20, font_color=WHITE, bold=True)

    facts = [
        ('1', '73% of COBOL systems depend on retiring experts.', 'The knowledge crisis is real.', RED),
        ('2', 'Domain training is required, not prompt engineering.', 'Generic AI cannot solve it.', ORANGE),
        ('3', '19 products, 245 manuals, and the only RAFT-trained platform.', 'TmaxSoft has the position.', GREEN),
    ]

    for i, (num, detail, label, color) in enumerate(facts):
        y = Inches(2.2 + i * 0.85)
        add_shape(slide, Inches(1.5), y, Inches(10), Inches(0.7), BG_CARD, BORDER)
        add_accent_bar(slide, Inches(1.5), y, Inches(0.7), color)
        add_text_box(slide, Inches(1.8), y + Pt(4), Inches(3), Inches(0.3),
                     label, font_size=12, font_color=color, bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.3), Inches(9.5), Inches(0.3),
                     detail, font_size=12, font_color=TEXT_SECONDARY)

    add_divider_line(slide, Inches(1.5), Inches(5.0), Inches(10))

    add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(1.0),
                 '"In five years, every legacy modernization project\n'
                 'will require AI intelligence.\n'
                 'The question is not whether \u2014 it is whose AI.\n'
                 'We intend it to be ours."',
                 font_size=18, font_color=WHITE, bold=False,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
                 'TmaxSoft Japan  \u00b7  AI KMS Platform Division',
                 font_size=12, font_color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 24, TOTAL_SLIDES)


def slide_25_next_steps(prs):
    """Next Steps & Contact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4),
                 'NEXT STEPS', font_size=10, font_color=TEXT_DIM, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
                 'Ready to Demonstrate, Deploy, and Deliver', font_size=32, font_color=WHITE, bold=True)

    actions = [
        ('Executive Demo Session', 'Live AI KMS walkthrough \u2014 60 minutes', BLUE),
        ('Pilot Program Discussion', '3-month proof of value with selected customer', GREEN),
        ('Technical Deep-Dive', 'Architecture review for IT leadership', PURPLE),
        ('Business Case Development', 'Custom ROI model for your organization', ORANGE),
    ]

    for i, (title, desc, color) in enumerate(actions):
        y = Inches(2.0 + i * 1.1)
        add_shape(slide, Inches(0.8), y, Inches(7), Inches(0.9), BG_CARD, BORDER)
        add_accent_bar(slide, Inches(0.8), y, Inches(0.9), color)
        add_text_box(slide, Inches(1.2), y + Pt(6), Inches(6.2), Inches(0.35),
                     title, font_size=14, font_color=color, bold=True)
        add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(6.2), Inches(0.35),
                     desc, font_size=12, font_color=TEXT_SECONDARY)

    # Resources panel
    add_shape(slide, Inches(8.3), Inches(2.0), Inches(4.2), Inches(4.3), BG_CARD, BORDER)
    add_accent_bar(slide, Inches(8.3), Inches(2.0), Inches(4.3), CYAN)
    add_text_box(slide, Inches(8.6), Inches(2.15), Inches(3.7), Inches(0.35),
                 'Available Resources', font_size=13, font_color=CYAN, bold=True)

    resources = [
        'Technical Architecture Document',
        'Customer ROI Calculator',
        'Live Demo Environment Access',
        'RAFT Training Methodology Paper',
        'Product Coverage Matrix (19 products)',
        'E2E Quality Test Results',
    ]
    add_bullet_list(slide, Inches(8.6), Inches(2.6), Inches(3.7), Inches(3.5),
                    resources, font_size=11, font_color=TEXT_SECONDARY, spacing=Pt(10))

    # Contact
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
                 'TmaxSoft Japan  |  AI Platform Division  |  openframe-ai-kms@tmaxsoft.com',
                 font_size=12, font_color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 25, TOTAL_SLIDES)


# ============================================================
# MAIN
# ============================================================
def generate():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_knowledge_fragmentation(prs)
    slide_04_why_generic_fails(prs)
    slide_05_why_now(prs)
    slide_06_product_overview(prs)
    slide_07_architecture(prs)
    slide_08_comparison(prs)
    slide_09_business_impact(prs)
    slide_10_roi(prs)
    slide_11_capability_support(prs)
    slide_12_capability_log(prs)
    slide_13_capability_migration(prs)
    slide_14_capability_docs(prs)
    slide_15_premium_support(prs)
    slide_16_strategic_position(prs)
    slide_17_business_model(prs)
    slide_18_competitive_moat(prs)
    slide_19_why_not_generic(prs)
    slide_20_raft(prs)
    slide_21_future_vision(prs)
    slide_22_global_expansion(prs)
    slide_23_investment(prs)
    slide_24_closing(prs)
    slide_25_next_steps(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, '..', 'docs', 'TmaxSoft_AI_KMS_Executive_Deck.pptx')
    output_path = os.path.normpath(output_path)

    prs.save(output_path)
    print(f'Generated: {output_path}')
    print(f'Slides: {len(prs.slides)}')
    return output_path


if __name__ == '__main__':
    generate()
