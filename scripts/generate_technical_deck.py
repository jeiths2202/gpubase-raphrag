#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenFrame AI KMS 기술 구현 보고서 - TmaxSoft 템플릿 기반 PowerPoint 생성기.
기술적 구현 중심, 한국어, 검증되지 않은 수치 최소화.
"""

import sys
import io
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

TEMPLATE_PATH = "TmaxSoft-ppt-template.pptx"
OUTPUT_PATH = "docs/TmaxSoft_AI_KMS_Technical_Deck.pptx"


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def clone_slide(prs, index):
    """Clone a slide preserving all shapes, images, and styling."""
    src = prs.slides[index]
    new_slide = prs.slides.add_slide(src.slide_layout)

    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = etree.QName(child.tag).localname
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            sp_tree.remove(child)

    rid_map = {}
    for rel in src.part.rels.values():
        if rel.reltype == RT.SLIDE_LAYOUT:
            for nr in new_slide.part.rels.values():
                if nr.reltype == RT.SLIDE_LAYOUT:
                    rid_map[rel.rId] = nr.rId
                    break
            continue
        try:
            if rel.is_external:
                new_rid = new_slide.part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref
                )
            else:
                new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rid
        except Exception:
            pass

    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for child in src.shapes._spTree:
        tag = etree.QName(child.tag).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        nc = deepcopy(child)
        for elem in nc.iter():
            for ak in list(elem.attrib.keys()):
                if R_NS in ak:
                    old = elem.get(ak)
                    if old in rid_map:
                        elem.set(ak, rid_map[old])
        sp_tree.append(nc)

    return new_slide


def delete_slide(prs, index):
    """Delete slide at given index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def set_text(slide, shape_name, new_text):
    """Find shape by name on a specific slide and replace its text."""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run._r.getparent().remove(run._r)
                    return True
    return False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title(prs):
    """Slide 1: 표지."""
    s = clone_slide(prs, 0)
    set_text(s, "Text 2", "© 2026. TmaxSoft. All Rights Reserved.")
    set_text(s, "Text 4", "OpenFrame")
    set_text(s, "Text 5", "AI KMS")
    set_text(s, "Text 6", "OpenFrame AI KMS")
    set_text(s, "Text 7", "기술 구현 보고서")
    set_text(s, "Text 8", "Agentic RAG + QLoRA 도메인 특화 학습 기반")
    set_text(s, "Text 9", "차세대 AI 지식관리 플랫폼")
    set_text(s, "Text 10", "2026년 3월")
    set_text(s, "Text 11", "작성자: 신이재")
    return s


def build_toc(prs):
    """Slide 2: 목차."""
    s = clone_slide(prs, 1)
    items = {
        "Text 17": "프로젝트 개요",
        "Text 19": "시스템 아키텍처",
        "Text 21": "핵심 기능",
        "Text 23": "Agentic RAG 파이프라인",
        "Text 25": "QLoRA 도메인 특화 학습",
        "Text 27": "Multi-Agent 협업 패턴",
        "Text 29": "지식 그래프 & 하이브리드 검색",
        "Text 31": "향후 로드맵",
        "Text 33": "감사합니다",
        # Clear unused items
        "Text 34": "",
        "Text 35": "",
        "Text 36": "",
        "Text 37": "",
    }
    for name, text in items.items():
        set_text(s, name, text)
    return s


def build_content(prs, badge, page, title,
                  left1, left2, left_body,
                  h1, b1, h2, b2, h3, bul1, bul2):
    """Build a content slide (Slide 3 pattern: left panel + right 3 sections)."""
    s = clone_slide(prs, 2)
    set_text(s, "Text 22", badge)
    set_text(s, "Text 20", title)
    set_text(s, "Text 23", str(page))
    # Left panel
    set_text(s, "Text 16", left1)
    set_text(s, "Text 17", left2)
    set_text(s, "Text 18", left_body)
    # Right sections
    set_text(s, "Text 6", h1)
    set_text(s, "Text 7", b1)
    set_text(s, "Text 8", h2)
    set_text(s, "Text 9", b2)
    set_text(s, "Text 10", h3)
    set_text(s, "Text 11", bul1)
    set_text(s, "Text 12", bul2)
    set_text(s, "Text 13", "")  # overflow area
    set_text(s, "Text 14", "")
    return s


def build_cards(prs, badge, page, title, cards):
    """Build a 6-card feature grid (Slide 4 pattern).

    cards: list of 6 dicts {title, badge, b1, b2, b3}.
    """
    s = clone_slide(prs, 3)
    set_text(s, "Text 58", badge)
    set_text(s, "Text 56", title)
    set_text(s, "Text 59", str(page))

    title_names = ["Text 6", "Text 11", "Text 16",
                   "Text 21", "Text 26", "Text 31"]
    badge_names = ["Text 32", "Text 33", "Text 34",
                   "Text 35", "Text 36", "Text 37"]
    bullet_groups = [
        ["Text 38", "Text 39", "Text 40"],
        ["Text 41", "Text 42", "Text 43"],
        ["Text 44", "Text 45", "Text 46"],
        ["Text 47", "Text 48", "Text 49"],
        ["Text 50", "Text 51", "Text 52"],
        ["Text 53", "Text 54", "Text 55"],
    ]

    for i, c in enumerate(cards):
        set_text(s, title_names[i], c["title"])
        set_text(s, badge_names[i], c["badge"])
        set_text(s, bullet_groups[i][0], c["b1"])
        set_text(s, bullet_groups[i][1], c["b2"])
        set_text(s, bullet_groups[i][2], c["b3"])
    return s


def build_closing(prs):
    """Last slide: 감사합니다."""
    return clone_slide(prs, 4)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation(TEMPLATE_PATH)

    # ── 1. 표지 ──────────────────────────────────────────────
    build_title(prs)
    print("  1/11  표지")

    # ── 2. 목차 ──────────────────────────────────────────────
    build_toc(prs)
    print("  2/11  목차")

    # ── 3. 프로젝트 개요 ─────────────────────────────────────
    build_content(
        prs, "01", 3, "프로젝트 개요",
        "OpenFrame 기술 지식의",
        "AI 기반 통합 관리",
        "TmaxSoft 제품군의 기술 매뉴얼을 AI가 학습하여 "
        "자연어 기반의 정확한 기술 지원을 제공하는 "
        "차세대 지식관리시스템입니다.",
        # Section 1
        "프로젝트 목적",
        "TmaxSoft OpenFrame 제품군의 방대한 기술 매뉴얼을 "
        "AI가 학습·검색하여, 다국어(한/일/영) 자연어 질의에 "
        "정확한 기술 지원을 제공하는 차세대 KMS 구축. "
        "제품별 전용 AI Agent가 자율적으로 검색 전략을 "
        "수립하고 근거 기반 답변을 생성합니다.",
        # Section 2
        "핵심 배경",
        "기존 키워드 검색의 한계로 제품 간 교차 참조가 어렵고, "
        "도메인 전문 지식이 분산되어 효율적 기술 지원이 불가능. "
        "범용 LLM은 TmaxSoft 제품 도메인 지식이 부재하여 "
        "환각(Hallucination) 발생. QLoRA 도메인 특화 학습으로 해결.",
        # Section 3
        "주요 목표",
        "제품별 QLoRA 어댑터 기반 도메인 특화 AI 응답 생성",
        "Graph + Vector + BM25 하이브리드 검색으로 정밀도 향상",
    )
    print("  3/11  프로젝트 개요")

    # ── 4. 시스템 아키텍처 ────────────────────────────────────
    build_content(
        prs, "02", 4, "시스템 아키텍처",
        "GPU 기반",
        "분산 AI 아키텍처",
        "FastAPI + Neo4j + vLLM 기반의 고성능 비동기 "
        "파이프라인으로, 제품별 전용 Agent가 동적으로 "
        "생성되어 자율적으로 검색과 응답을 수행합니다.",
        # Section 1
        "Backend 구성",
        "FastAPI (Python 3.10+) 비동기 SSE 스트리밍 기반. "
        "Router → Service → Repository 레이어 분리 아키텍처. "
        "ManualRegistryService가 PDF 디렉토리에서 제품을 "
        "동적으로 발견하여 제품별 전용 Agent를 자동 생성합니다.",
        # Section 2
        "데이터베이스 & 검색",
        "Neo4j Graph + Vector Index 이중 활용. "
        "Chunk-Entity MENTIONS 관계 모델링으로 지식 네트워크 구축. "
        "파일 시스템 기반 Summary 사전(에러코드, 명령어, 용어)으로 "
        "쿼리 보강 후 검색 정밀도를 향상시킵니다.",
        # Section 3
        "LLM 엔진",
        "vLLM 기반 Qwen 2.5 + 제품별 QLoRA 어댑터 동적 로딩",
        "NVIDIA A100 GPU 클러스터, Tensor Parallel 추론",
    )
    print("  4/11  시스템 아키텍처")

    # ── 5. 핵심 기능 (6-card grid) ───────────────────────────
    build_cards(
        prs, "03", 5, "핵심 기능 (Core Features)",
        [
            {
                "title": "Agentic RAG 검색",
                "badge": "Agent",
                "b1": "제품별 전용 Agent 자율 검색 전략 수립",
                "b2": "Two-Stage Retrieval: Summary → Vector/Graph",
                "b3": "ResponseVerifier 기반 문장별 환각 검증",
            },
            {
                "title": "QLoRA 도메인 학습",
                "badge": "Training",
                "b1": "3-Phase: CPT → SFT → DPO 파이프라인",
                "b2": "제품별 어댑터 vLLM 동적 로딩",
                "b3": "RAFT 방법론: Distractor 문서 구분 학습",
            },
            {
                "title": "지식 그래프 탐색",
                "badge": "Graph",
                "b1": "엔티티 자동 추출 및 관계 모델링",
                "b2": "MENTIONS 기반 관련 Chunk 연결 탐색",
                "b3": "D3.js 인터랙티브 마인드맵 시각화",
            },
            {
                "title": "Multi-Agent 협업",
                "badge": "Teams",
                "b1": "병렬 검색, 경쟁 가설 등 5가지 협업 패턴",
                "b2": "asyncio.gather 기반 동시 처리",
                "b3": "Feature flag 기반 점진적 활성화",
            },
            {
                "title": "웹 문서 고속 검색",
                "badge": "WebDoc",
                "b1": "docs.tmaxsoft.com 실시간 크롤링 인덱싱",
                "b2": "IDF 기반 고속 키워드 매칭 (<10ms)",
                "b3": "고신뢰 매칭 시 PDF RAG 파이프라인 우회",
            },
            {
                "title": "환각 방지 시스템",
                "badge": "Verify",
                "b1": "문장별 코사인 유사도 검증 필터링",
                "b2": "E2E Hallucination 테스트 자동화",
                "b3": "RAFT: Oracle vs Distractor 문서 구분",
            },
        ],
    )
    print("  5/11  핵심 기능")

    # ── 6. Agentic RAG 파이프라인 ─────────────────────────────
    build_content(
        prs, "04", 6, "Agentic RAG 파이프라인",
        "제품별 전문 Agent",
        "자율 검색 파이프라인",
        "QueryRouter가 자연어 질문을 분석하여 적합한 "
        "제품 Agent로 자동 라우팅합니다. Agent가 자율적으로 "
        "검색 전략을 수립하고, 근거 문서 기반으로만 "
        "응답을 생성합니다.",
        # Section 1
        "QueryRouter (LLM 미사용)",
        "자연어 질문에서 키워드+패턴 기반으로 제품 자동 라우팅. "
        "확정(conf>=0.8), 후보(0.5~0.8), 미매칭(<0.5) 3단계 분류. "
        "ProductContextMemory를 통한 세션 간 라우팅 컨텍스트 영속화. "
        "검색 단계에서 LLM 미사용으로 결정론적 동작 보장.",
        # Section 2
        "Two-Stage Retrieval",
        "Stage 1: Summary 사전(에러코드/명령어/용어)에서 "
        "컨텍스트 보강 (<10ms). "
        "Stage 2: 보강된 쿼리로 Vector + Graph DB 하이브리드 검색 수행. "
        "검색 결과 기반으로만 LLM 응답을 생성하여 환각을 최소화.",
        # Section 3
        "ResponseVerifier",
        "LLM 생성 답변의 문장별 코사인 유사도 검증",
        "근거 없는 문장 자동 필터링 및 경고 표시",
    )
    print("  6/11  Agentic RAG")

    # ── 7. QLoRA 도메인 특화 학습 ─────────────────────────────
    build_content(
        prs, "05", 7, "QLoRA 도메인 특화 학습",
        "3-Phase",
        "학습 파이프라인",
        "도메인 지식 주입(CPT)부터 명령 수행 학습(SFT), "
        "선호도 정렬(DPO)까지 3단계 파이프라인으로 "
        "제품별 전문 AI를 구축합니다. PDCA 반복 정제로 "
        "데이터 품질을 지속 개선합니다.",
        # Section 1
        "Phase 1: CPT (도메인 지식 주입)",
        "PDF 원문 텍스트 기반 Continued Pre-Training. "
        "Plain Text 포맷으로 4096 토큰 청크 학습. "
        "LoRA r=64, α=128 설정. "
        "<|endoftext|> 토큰으로 문서 경계 구분. "
        "모델이 OpenFrame 도메인 용어와 개념을 자연스럽게 이해.",
        # Section 2
        "Phase 2-3: SFT → DPO",
        "SFT: 제품별 Q-A 쌍 ChatML 포맷으로 명령 수행 학습. "
        "제품별 QLoRA 어댑터 생성 후 vLLM 동적 로딩. "
        "DPO: 교차 제품 Distractor 학습으로 환각 억제. "
        "RAFT 방법론 기반 Oracle/Distractor 문서 구분 능력 강화.",
        # Section 3
        "데이터셋 PDCA 정제",
        "v4→v9 반복 정제: 시맨틱 클리닝, 패턴 필터링",
        "코사인 유사도 0.95+ 중복 제거, Q-A 불일치 필터링",
    )
    print("  7/11  QLoRA")

    # ── 8. Multi-Agent 협업 패턴 (6-card grid) ────────────────
    build_cards(
        prs, "06", 8, "Multi-Agent 협업 패턴",
        [
            {
                "title": "병렬 검색 (Pattern A)",
                "badge": "Parallel",
                "b1": "Web Doc + PDF RAG 동시 검색",
                "b2": "asyncio.gather 기반 비동기 병렬 처리",
                "b3": "결과 병합 후 최적 소스 자동 선택",
            },
            {
                "title": "경쟁 가설 (Pattern B)",
                "badge": "Compete",
                "b1": "다중 temperature 기반 가설 생성",
                "b2": "룰 기반 평가 → 최적 답변 선택",
                "b3": "답변 품질 다양성 확보",
            },
            {
                "title": "멀티 LoRA (Pattern C)",
                "badge": "LoRA",
                "b1": "관련 제품 QLoRA 어댑터 병렬 로딩",
                "b2": "confidence 기반 최적 어댑터 선택",
                "b3": "제품 경계가 모호한 질문 대응",
            },
            {
                "title": "제품 비교 (Pattern D)",
                "badge": "Compare",
                "b1": "제품별 독립 검색 + LLM 병렬 실행",
                "b2": "비교 합성 답변 자동 생성",
                "b3": "멀티 제품 교차 분석 지원",
            },
            {
                "title": "자기 개선 (Pattern E)",
                "badge": "Improve",
                "b1": "사용자 피드백 JSONL 축적",
                "b2": "QLoRA 증분 재학습 데이터 활용",
                "b3": "지속적 품질 개선 루프 구축",
            },
            {
                "title": "TeamOrchestrator",
                "badge": "Core",
                "b1": "AgenticRAGService 래핑 (중앙 조율)",
                "b2": "Feature flag 기반 점진적 활성화",
                "b3": "Flag OFF 시 기존 동작 100% 유지",
            },
        ],
    )
    print("  8/11  Multi-Agent")

    # ── 9. 지식 그래프 & 하이브리드 검색 ──────────────────────
    build_content(
        prs, "07", 9, "지식 그래프 & 하이브리드 검색",
        "3중 검색",
        "파이프라인",
        "Vector + Graph + BM25 검색을 결합하고, "
        "Summary 사전으로 쿼리를 보강하여 검색 정밀도를 "
        "극대화합니다. 엔티티 자동 추출로 지식 네트워크를 "
        "구축합니다.",
        # Section 1
        "Chunk-Entity 파이프라인",
        "Neo4j Graph DB 기반 Chunk-Entity 관계 모델링. "
        "3-Phase 엔티티 추출: Summary 사전 매칭(conf=0.95) → "
        "정규식 패턴(0.80) → 카타카나 폴백(0.70). "
        "엔티티 유형: config, command, concept, error_code, product. "
        "MERGE 기반 멱등성으로 안전한 재실행 보장.",
        # Section 2
        "하이브리드 검색 전략",
        "Vector: NV-EmbedQA-Mistral 임베딩 코사인 유사도 검색. "
        "Graph: Neo4j Cypher 쿼리로 엔티티 관계 기반 연관 문서 탐색. "
        "BM25: 에러 코드, 명령어 등 정확한 키워드 매칭. "
        "Web Doc Fast Path: docs.tmaxsoft.com IDF 기반 고속 검색.",
        # Section 3
        "Summary Two-Stage 보강",
        "에러코드/명령어/용어 사전에서 컨텍스트 보강 (<10ms)",
        "보강된 쿼리로 Vector + Graph DB 정밀 검색 수행",
    )
    print("  9/11  지식 그래프")

    # ── 10. 향후 로드맵 ──────────────────────────────────────
    build_content(
        prs, "08", 10, "향후 로드맵",
        "확장 및",
        "고도화 계획",
        "모델 업그레이드, 메모리 영속화, 자동 재학습 등 "
        "지속적인 시스템 고도화를 통해 AI 기술 지원 "
        "플랫폼으로 발전시킵니다.",
        # Section 1
        "모델 & 인프라 업그레이드",
        "Qwen 32B 모델 업그레이드 (32K+ 컨텍스트 윈도우). "
        "A100 48GB 텐서 병렬 배포로 장문 문서 처리 능력 강화. "
        "vLLM continuous batching으로 동시 다수 사용자 "
        "처리 최적화.",
        # Section 2
        "시스템 고도화",
        "PostgresStore 기반 장기 메모리 영속화 "
        "(현재 InMemoryStore 개발용). "
        "자동 재학습 파이프라인: 사용자 피드백 JSONL → "
        "QLoRA 증분 학습. "
        "JCL 진단 보고서 자동 생성 (HTML 템플릿, 5-Agent 파이프라인).",
        # Section 3
        "글로벌 확장",
        "일본어/한국어/영어 동시 지원 강화",
        "Fujitsu XSP/AIM 메인프레임 마이그레이션 지원 확대",
    )
    print(" 10/11  향후 로드맵")

    # ── 11. 감사합니다 ───────────────────────────────────────
    build_closing(prs)
    print(" 11/11  감사합니다")

    # ── Delete original 5 template slides ────────────────────
    for _ in range(5):
        delete_slide(prs, 0)

    prs.save(OUTPUT_PATH)
    print(f"\nGenerated: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
