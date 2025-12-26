"""
Document Analyzer Service

Analyzes documents to determine visual characteristics and complexity.
Used for routing decisions between Vision LLM and Text LLM.
"""

import io
import logging
import re
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from PIL import Image

from app.api.models.vision import (
    DocumentVisualProfile,
    ImageType,
    ProcessingMode,
)

logger = logging.getLogger(__name__)


# MIME type mappings
PURE_IMAGE_MIMES = {
    "image/png", "image/jpeg", "image/jpg", "image/gif",
    "image/bmp", "image/tiff", "image/webp"
}

DOCUMENT_MIMES = {
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

TEXT_MIMES = {
    "text/plain", "text/markdown", "text/csv", "text/html",
    "application/json", "application/xml"
}

EXTENSION_TO_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".tiff": "image/tiff",
    ".webp": "image/webp",
    ".pdf": "application/pdf",
    ".doc": "application/msword",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xls": "application/vnd.ms-excel",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".ppt": "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".html": "text/html",
    ".json": "application/json",
}


@dataclass
class ChartDetectionResult:
    """Result of chart/graph detection in an image"""
    has_chart: bool
    chart_type: Optional[str] = None  # "bar", "line", "pie", etc.
    confidence: float = 0.0


@dataclass
class DiagramDetectionResult:
    """Result of diagram detection in an image"""
    has_diagram: bool
    diagram_type: Optional[str] = None  # "flowchart", "sequence", etc.
    confidence: float = 0.0


class DocumentAnalyzer:
    """
    Analyzes document visual characteristics.

    Three-layer analysis:
    1. MIME Type & Extension - Basic classification
    2. Content Analysis - Image/text ratio, density
    3. Visual Complexity - Charts, diagrams, screenshots

    Usage:
        analyzer = DocumentAnalyzer()
        profile = await analyzer.analyze(file_path, mime_type)

        if profile.requires_vision_llm:
            # Route to Vision LLM
            ...
    """

    # Thresholds for visual complexity detection
    IMAGE_AREA_THRESHOLD = 0.30  # 30% of page area
    TEXT_DENSITY_THRESHOLD = 100  # chars per page
    EXTRACTABLE_TEXT_THRESHOLD = 0.50  # 50% of expected text

    # Visual complexity weights
    WEIGHTS = {
        "image_area_ratio": 0.25,
        "text_density_inverse": 0.20,
        "has_charts": 0.20,
        "has_diagrams": 0.15,
        "has_screenshots": 0.10,
        "requires_ocr": 0.10,
    }

    def __init__(
        self,
        image_area_threshold: float = 0.30,
        text_density_threshold: float = 100,
        enable_ml_detection: bool = False,
    ):
        """
        Initialize analyzer.

        Args:
            image_area_threshold: Threshold for image area ratio
            text_density_threshold: Minimum chars per page for text-heavy
            enable_ml_detection: Enable ML-based chart/diagram detection
        """
        self.image_area_threshold = image_area_threshold
        self.text_density_threshold = text_density_threshold
        self.enable_ml_detection = enable_ml_detection

    async def analyze(
        self,
        file_path: Path,
        mime_type: Optional[str] = None,
        document_id: Optional[str] = None,
    ) -> DocumentVisualProfile:
        """
        Analyze document and create visual profile.

        Args:
            file_path: Path to the document
            mime_type: Optional MIME type (auto-detected if not provided)
            document_id: Optional document ID for the profile

        Returns:
            DocumentVisualProfile with all metrics
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        # Determine MIME type
        if not mime_type:
            mime_type = EXTENSION_TO_MIME.get(extension, "application/octet-stream")

        # Create base profile
        profile = DocumentVisualProfile(
            document_id=document_id or str(file_path),
            mime_type=mime_type,
            extension=extension,
            is_pure_image=mime_type in PURE_IMAGE_MIMES,
        )

        # Layer 1: Basic classification based on MIME type
        if profile.is_pure_image:
            profile = await self._analyze_image(file_path, profile)
        elif mime_type == "application/pdf":
            profile = await self._analyze_pdf(file_path, profile)
        elif mime_type in DOCUMENT_MIMES:
            profile = await self._analyze_office(file_path, profile)
        elif mime_type in TEXT_MIMES:
            profile = await self._analyze_text(file_path, profile)

        # Calculate visual complexity score
        profile.visual_complexity_score = self._calculate_complexity_score(profile)

        return profile

    async def analyze_bytes(
        self,
        content: bytes,
        mime_type: str,
        document_id: Optional[str] = None,
    ) -> DocumentVisualProfile:
        """
        Analyze document from bytes.

        Args:
            content: Document content as bytes
            mime_type: MIME type of the content
            document_id: Optional document ID

        Returns:
            DocumentVisualProfile
        """
        extension = ""
        for ext, mt in EXTENSION_TO_MIME.items():
            if mt == mime_type:
                extension = ext
                break

        profile = DocumentVisualProfile(
            document_id=document_id or "bytes_document",
            mime_type=mime_type,
            extension=extension,
            is_pure_image=mime_type in PURE_IMAGE_MIMES,
        )

        if profile.is_pure_image:
            profile = await self._analyze_image_bytes(content, profile)
        elif mime_type == "application/pdf":
            profile = await self._analyze_pdf_bytes(content, profile)

        profile.visual_complexity_score = self._calculate_complexity_score(profile)
        return profile

    async def _analyze_image(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze a pure image file."""
        try:
            with open(file_path, "rb") as f:
                image_bytes = f.read()
            return await self._analyze_image_bytes(image_bytes, profile)
        except Exception as e:
            logger.error(f"Error analyzing image {file_path}: {e}")
            return profile

    async def _analyze_image_bytes(
        self,
        image_bytes: bytes,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze image from bytes."""
        try:
            image = Image.open(io.BytesIO(image_bytes))

            profile.total_pages = 1
            profile.image_count = 1
            profile.image_area_ratio = 1.0  # Pure image
            profile.text_density = 0.0

            # Detect visual elements
            chart_result = await self._detect_chart(image_bytes)
            diagram_result = await self._detect_diagram(image_bytes)
            screenshot_result = self._detect_screenshot(image)

            profile.has_charts = chart_result.has_chart
            profile.has_diagrams = diagram_result.has_diagram
            profile.has_screenshots = screenshot_result

            # Check if OCR is needed (no extractable text from pure image)
            profile.requires_ocr = True
            profile.extractable_text_ratio = 0.0

        except Exception as e:
            logger.error(f"Error analyzing image bytes: {e}")

        return profile

    async def _analyze_pdf(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze a PDF document."""
        try:
            with open(file_path, "rb") as f:
                pdf_bytes = f.read()
            return await self._analyze_pdf_bytes(pdf_bytes, profile)
        except Exception as e:
            logger.error(f"Error analyzing PDF {file_path}: {e}")
            return profile

    async def _analyze_pdf_bytes(
        self,
        pdf_bytes: bytes,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze PDF from bytes."""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            profile.total_pages = len(doc)

            total_text_chars = 0
            total_image_area = 0.0
            total_page_area = 0.0
            image_count = 0
            table_count = 0

            for page_num in range(len(doc)):
                page = doc[page_num]
                page_rect = page.rect
                page_area = page_rect.width * page_rect.height
                total_page_area += page_area

                # Extract text
                text = page.get_text()
                total_text_chars += len(text)

                # Count and measure images
                images = page.get_images(full=True)
                for img in images:
                    image_count += 1
                    # Estimate image area (simplified)
                    # Full implementation would get actual image bounds
                    total_image_area += page_area * 0.2  # Conservative estimate

                # Detect tables (simplified: look for table-like text patterns)
                if self._has_table_pattern(text):
                    table_count += 1

            doc.close()

            # Calculate metrics
            profile.image_count = image_count
            profile.table_count = table_count

            if total_page_area > 0:
                profile.image_area_ratio = min(1.0, total_image_area / total_page_area)
                profile.table_area_ratio = min(0.5, table_count * 0.1)  # Estimate

            if profile.total_pages > 0:
                profile.text_density = total_text_chars / profile.total_pages

            # Determine if OCR is needed
            expected_chars = profile.total_pages * 500  # ~500 chars per page expected
            if expected_chars > 0:
                profile.extractable_text_ratio = min(1.0, total_text_chars / expected_chars)
            profile.requires_ocr = profile.extractable_text_ratio < self.EXTRACTABLE_TEXT_THRESHOLD

            # Visual element detection (sample first few pages)
            await self._detect_visual_elements_in_pdf(pdf_bytes, profile)

        except ImportError:
            logger.warning("PyMuPDF not available for PDF analysis")
        except Exception as e:
            logger.error(f"Error analyzing PDF bytes: {e}")

        return profile

    async def _analyze_office(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze Office document (Word, Excel, PowerPoint)."""
        try:
            extension = file_path.suffix.lower()

            if extension in (".pptx", ".ppt"):
                return await self._analyze_powerpoint(file_path, profile)
            elif extension in (".docx", ".doc"):
                return await self._analyze_word(file_path, profile)
            elif extension in (".xlsx", ".xls"):
                return await self._analyze_excel(file_path, profile)

        except Exception as e:
            logger.error(f"Error analyzing Office document {file_path}: {e}")

        return profile

    async def _analyze_powerpoint(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze PowerPoint presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            profile.total_pages = len(prs.slides)

            image_count = 0
            total_text = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        total_text += len(shape.text)
                    if hasattr(shape, "image"):
                        image_count += 1

            profile.image_count = image_count
            profile.text_density = total_text / max(1, profile.total_pages)

            # PowerPoint is typically visual-heavy
            profile.image_area_ratio = min(1.0, image_count * 0.15)

        except ImportError:
            logger.warning("python-pptx not available")
        except Exception as e:
            logger.error(f"Error analyzing PowerPoint: {e}")

        return profile

    async def _analyze_word(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze Word document."""
        try:
            from docx import Document

            doc = Document(str(file_path))
            profile.total_pages = 1  # Word doesn't expose page count easily

            total_text = 0
            image_count = 0
            table_count = len(doc.tables)

            for para in doc.paragraphs:
                total_text += len(para.text)

            # Count inline shapes (images)
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    image_count += 1

            profile.image_count = image_count
            profile.table_count = table_count
            profile.text_density = total_text
            profile.image_area_ratio = min(0.5, image_count * 0.1)
            profile.table_area_ratio = min(0.5, table_count * 0.1)

        except ImportError:
            logger.warning("python-docx not available")
        except Exception as e:
            logger.error(f"Error analyzing Word document: {e}")

        return profile

    async def _analyze_excel(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze Excel spreadsheet."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), data_only=True)
            profile.total_pages = len(wb.sheetnames)

            total_cells = 0
            chart_count = 0

            for sheet in wb:
                total_cells += sheet.max_row * sheet.max_column
                chart_count += len(sheet._charts)

            profile.text_density = total_cells / max(1, profile.total_pages)
            profile.has_charts = chart_count > 0
            profile.table_area_ratio = 0.8  # Excel is essentially all tables

        except ImportError:
            logger.warning("openpyxl not available")
        except Exception as e:
            logger.error(f"Error analyzing Excel: {e}")

        return profile

    async def _analyze_text(
        self,
        file_path: Path,
        profile: DocumentVisualProfile
    ) -> DocumentVisualProfile:
        """Analyze text-based document."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            profile.total_pages = 1
            profile.text_density = len(content)
            profile.extractable_text_ratio = 1.0
            profile.image_area_ratio = 0.0

        except Exception as e:
            logger.error(f"Error analyzing text file: {e}")

        return profile

    async def _detect_visual_elements_in_pdf(
        self,
        pdf_bytes: bytes,
        profile: DocumentVisualProfile
    ) -> None:
        """Detect charts, diagrams in PDF pages."""
        try:
            import fitz

            doc = fitz.open(stream=pdf_bytes, filetype="pdf")

            # Sample first 3 pages for visual detection
            pages_to_check = min(3, len(doc))

            for page_num in range(pages_to_check):
                page = doc[page_num]

                # Render page to image for analysis
                mat = fitz.Matrix(1.0, 1.0)  # Low resolution for speed
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")

                # Detect visual elements
                chart_result = await self._detect_chart(img_bytes)
                diagram_result = await self._detect_diagram(img_bytes)

                if chart_result.has_chart:
                    profile.has_charts = True
                if diagram_result.has_diagram:
                    profile.has_diagrams = True

                # Early exit if both found
                if profile.has_charts and profile.has_diagrams:
                    break

            doc.close()

        except Exception as e:
            logger.warning(f"Visual element detection failed: {e}")

    async def _detect_chart(self, image_bytes: bytes) -> ChartDetectionResult:
        """
        Detect if image contains a chart/graph.

        Uses heuristics:
        - Color histogram analysis (charts have distinct color bands)
        - Edge detection patterns
        - Aspect ratio analysis
        """
        if not self.enable_ml_detection:
            # Simple heuristic-based detection
            return self._detect_chart_heuristic(image_bytes)

        # ML-based detection would go here
        return ChartDetectionResult(has_chart=False)

    def _detect_chart_heuristic(self, image_bytes: bytes) -> ChartDetectionResult:
        """Heuristic-based chart detection."""
        try:
            image = Image.open(io.BytesIO(image_bytes))

            # Convert to RGB if needed
            if image.mode != "RGB":
                image = image.convert("RGB")

            # Analyze color distribution
            colors = image.getcolors(maxcolors=256)
            if colors:
                # Charts typically have fewer, more distinct colors
                if len(colors) < 50:
                    # Check for typical chart color patterns
                    dominant_colors = sorted(colors, key=lambda x: x[0], reverse=True)[:10]

                    # Charts often have blue, red, green bars
                    chart_colors = 0
                    for count, color in dominant_colors:
                        r, g, b = color
                        # Check for saturated colors (typical in charts)
                        saturation = max(r, g, b) - min(r, g, b)
                        if saturation > 100:
                            chart_colors += 1

                    if chart_colors >= 2:
                        return ChartDetectionResult(
                            has_chart=True,
                            confidence=0.6
                        )

        except Exception as e:
            logger.debug(f"Chart detection error: {e}")

        return ChartDetectionResult(has_chart=False)

    async def _detect_diagram(self, image_bytes: bytes) -> DiagramDetectionResult:
        """
        Detect if image contains a diagram.

        Uses heuristics:
        - Line and arrow detection
        - Box/shape detection
        - Connected component analysis
        """
        if not self.enable_ml_detection:
            return self._detect_diagram_heuristic(image_bytes)

        return DiagramDetectionResult(has_diagram=False)

    def _detect_diagram_heuristic(self, image_bytes: bytes) -> DiagramDetectionResult:
        """Heuristic-based diagram detection."""
        try:
            image = Image.open(io.BytesIO(image_bytes))

            # Simple heuristic: diagrams often have high contrast
            # and geometric shapes
            if image.mode != "L":
                image = image.convert("L")

            # Check for high contrast (many pure black/white pixels)
            histogram = image.histogram()
            dark_pixels = sum(histogram[:30])
            light_pixels = sum(histogram[225:])
            total_pixels = sum(histogram)

            contrast_ratio = (dark_pixels + light_pixels) / total_pixels

            if contrast_ratio > 0.3:
                return DiagramDetectionResult(
                    has_diagram=True,
                    confidence=0.5
                )

        except Exception as e:
            logger.debug(f"Diagram detection error: {e}")

        return DiagramDetectionResult(has_diagram=False)

    def _detect_screenshot(self, image: Image.Image) -> bool:
        """Detect if image is likely a screenshot."""
        try:
            # Screenshots typically have:
            # - Rectangular shape
            # - Standard screen dimensions
            # - UI-like patterns

            width, height = image.size

            # Common screen aspect ratios
            aspect_ratio = width / height
            common_ratios = [16/9, 16/10, 4/3, 21/9]

            for ratio in common_ratios:
                if abs(aspect_ratio - ratio) < 0.1:
                    return True

            # Check for common screen sizes
            common_widths = [1920, 1280, 1366, 1440, 2560, 3840]
            if width in common_widths:
                return True

        except Exception as e:
            logger.debug(f"Screenshot detection error: {e}")

        return False

    def _has_table_pattern(self, text: str) -> bool:
        """Check if text contains table-like patterns."""
        lines = text.split("\n")

        # Look for lines with consistent spacing/tabs
        tab_lines = sum(1 for line in lines if "\t" in line or "  " in line)

        # Tables typically have multiple lines with tabs/spacing
        if tab_lines >= 3:
            return True

        # Look for patterns like "| cell | cell |"
        pipe_pattern = re.compile(r"\|.*\|")
        pipe_lines = sum(1 for line in lines if pipe_pattern.search(line))

        return pipe_lines >= 2

    def _calculate_complexity_score(
        self,
        profile: DocumentVisualProfile
    ) -> float:
        """
        Calculate visual complexity score (0.0 - 1.0).

        Formula:
        score = (
            image_area_ratio * 0.25 +
            (1 - text_density_normalized) * 0.20 +
            has_charts * 0.20 +
            has_diagrams * 0.15 +
            has_screenshots * 0.10 +
            requires_ocr * 0.10
        )
        """
        # Normalize text density (inverse: lower density = higher complexity)
        text_density_normalized = min(1.0, profile.text_density / 1000)
        text_complexity = 1.0 - text_density_normalized

        score = (
            profile.image_area_ratio * self.WEIGHTS["image_area_ratio"] +
            text_complexity * self.WEIGHTS["text_density_inverse"] +
            (1.0 if profile.has_charts else 0.0) * self.WEIGHTS["has_charts"] +
            (1.0 if profile.has_diagrams else 0.0) * self.WEIGHTS["has_diagrams"] +
            (1.0 if profile.has_screenshots else 0.0) * self.WEIGHTS["has_screenshots"] +
            (1.0 if profile.requires_ocr else 0.0) * self.WEIGHTS["requires_ocr"]
        )

        return min(1.0, max(0.0, score))

    def get_recommended_mode(
        self,
        profile: DocumentVisualProfile
    ) -> ProcessingMode:
        """Get recommended processing mode based on profile."""
        return profile.recommended_processing_mode
