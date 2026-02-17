"""
Document Parser Service
Handles parsing of various document formats: PDF, Word, Excel, PowerPoint, Text, Images
"""
import asyncio
import base64
import io
import os
import re
import uuid
from abc import ABC, abstractmethod
from datetime import datetime, timezone
from typing import Optional, List, Dict, Any, Tuple, BinaryIO

# PowerPoint parsing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Excel parsing
try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

# Word parsing
try:
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

from ..models.document import (
    DocumentType,
    ProcessingMode,
    SUPPORTED_MIME_TYPES,
    EXTENSION_TO_MIME,
    VLMExtractionResult,
    ImageInfo,
    TableInfo,
)


class ParsedDocument:
    """Represents a parsed document with extracted content."""

    def __init__(
        self,
        document_type: DocumentType,
        filename: str,
        mime_type: str
    ):
        self.document_type = document_type
        self.filename = filename
        self.mime_type = mime_type
        self.text_content: str = ""
        self.pages: List[Dict[str, Any]] = []
        self.tables: List[TableInfo] = []
        self.images: List[ImageInfo] = []
        self.metadata: Dict[str, Any] = {}
        self.chunks: List[Dict[str, Any]] = []
        self.processing_info: Dict[str, Any] = {}

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        return {
            "document_type": self.document_type.value,
            "filename": self.filename,
            "mime_type": self.mime_type,
            "text_content": self.text_content,
            "pages": self.pages,
            "tables": [t.model_dump() for t in self.tables],
            "images": [i.model_dump() for i in self.images],
            "metadata": self.metadata,
            "chunks": self.chunks,
            "processing_info": self.processing_info
        }


class BaseDocumentParser(ABC):
    """Abstract base class for document parsers."""

    @abstractmethod
    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        """Parse document and extract content."""
        pass

    @abstractmethod
    def supports(self, mime_type: str) -> bool:
        """Check if parser supports the given MIME type."""
        pass


class TextParser(BaseDocumentParser):
    """Parser for plain text files."""

    SUPPORTED_TYPES = [
        "text/plain",
        "text/markdown",
        "text/html",
        "text/csv",
        "application/json"
    ]

    def supports(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        options = options or {}
        encoding = options.get("encoding", "utf-8")

        # Determine document type
        ext = os.path.splitext(filename)[1].lower()
        mime_type = EXTENSION_TO_MIME.get(ext, "text/plain")
        doc_type = SUPPORTED_MIME_TYPES.get(mime_type, DocumentType.TEXT)

        doc = ParsedDocument(doc_type, filename, mime_type)

        try:
            # Try to decode with specified encoding
            text = file_content.decode(encoding)
        except UnicodeDecodeError:
            # Fallback to different encodings
            for enc in ["utf-8", "cp949", "euc-kr", "latin-1"]:
                try:
                    text = file_content.decode(enc)
                    encoding = enc
                    break
                except UnicodeDecodeError:
                    continue
            else:
                text = file_content.decode("utf-8", errors="replace")

        doc.text_content = text
        doc.metadata = {
            "encoding": encoding,
            "line_count": len(text.splitlines()),
            "char_count": len(text),
            "word_count": len(text.split())
        }

        # Parse specific formats
        if mime_type == "text/csv":
            doc.tables = await self._parse_csv(text)
        elif mime_type == "application/json":
            doc.metadata["json_parsed"] = True
        elif mime_type == "text/markdown":
            doc.metadata["has_code_blocks"] = "```" in text
            doc.metadata["has_headers"] = bool(re.search(r"^#+\s", text, re.MULTILINE))

        doc.pages = [{"page_number": 1, "content": text}]
        doc.processing_info = {
            "parser": "TextParser",
            "parsed_at": datetime.now(timezone.utc).isoformat()
        }

        return doc

    async def _parse_csv(self, text: str) -> List[TableInfo]:
        """Parse CSV content into table."""
        lines = text.strip().split("\n")
        if not lines:
            return []

        # Simple CSV parsing (in production, use csv module)
        delimiter = "," if "," in lines[0] else "\t"
        rows = [line.split(delimiter) for line in lines]

        if len(rows) < 2:
            return []

        table = TableInfo(
            id=f"table_{uuid.uuid4().hex[:8]}",
            headers=rows[0],
            rows=rows[1:],
            caption="CSV Data",
            markdown=self._rows_to_markdown(rows[0], rows[1:])
        )
        return [table]

    def _rows_to_markdown(self, headers: List[str], rows: List[List[str]]) -> str:
        """Convert table rows to markdown."""
        md = "| " + " | ".join(headers) + " |\n"
        md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        for row in rows:
            md += "| " + " | ".join(row) + " |\n"
        return md


class PDFParser(BaseDocumentParser):
    """Parser for PDF documents using pypdf for real extraction."""

    SUPPORTED_TYPES = ["application/pdf"]

    def __init__(self, vlm_service=None):
        """
        Initialize PDF parser.

        Args:
            vlm_service: Optional VLM service for OCR and enhanced extraction
        """
        self.vlm_service = vlm_service

    def supports(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        """
        Parse PDF document using pypdf library for real statistics.
        """
        options = options or {}
        doc = ParsedDocument(DocumentType.PDF, filename, "application/pdf")

        # Use real PDF parsing with pypdf
        pdf_data = await self._extract_pdf_content(file_content, options)

        doc.text_content = pdf_data["text_content"]
        doc.pages = pdf_data["pages"]
        doc.metadata = pdf_data["metadata"]
        doc.images = pdf_data["images"]
        doc.tables = pdf_data["tables"]

        doc.processing_info = {
            "parser": "PDFParser",
            "parsed_at": datetime.now(timezone.utc).isoformat(),
            "options": options,
            "extraction_method": "pypdf"
        }

        return doc

    async def _extract_pdf_content(
        self,
        content: bytes,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Extract real content from PDF using pypdf, with optional VLM OCR."""
        processing_mode = options.get("processing_mode", "text_only")
        enable_vlm = options.get("enable_vlm", False)
        language = options.get("language", "auto")

        try:
            from pypdf import PdfReader

            # Run in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                None,
                self._sync_extract_pdf,
                content,
                options
            )

            # Check if VLM OCR is needed
            if (processing_mode in ["image_only", "vlm_enhanced"] or enable_vlm) and self.vlm_service:
                # Enhance pages with low text content using VLM OCR
                result = await self._enhance_with_vlm_ocr(
                    content, result, processing_mode, language
                )

            return result

        except ImportError:
            # Fallback to mock if pypdf not available
            return await self._fallback_mock_extraction(content)
        except Exception as e:
            # Fallback on any error
            print(f"PDF extraction error: {e}, using fallback")
            return await self._fallback_mock_extraction(content)

    async def _enhance_with_vlm_ocr(
        self,
        pdf_content: bytes,
        parsed_result: Dict[str, Any],
        processing_mode: str,
        language: str = "auto"
    ) -> Dict[str, Any]:
        """
        Enhance PDF extraction with VLM-based OCR for pages with little/no text.

        Args:
            pdf_content: Original PDF bytes
            parsed_result: Initial parsing result from pypdf
            processing_mode: Processing mode (ocr, vlm_enhanced, multimodal)
            language: Expected document language

        Returns:
            Enhanced parsing result with OCR text
        """
        if not self.vlm_service:
            return parsed_result

        try:
            from . import pdf_compat

            total_pages = pdf_compat.get_page_count(pdf_content)
            pages = parsed_result.get("pages", [])
            enhanced_pages = []
            full_text_parts = []
            ocr_count = 0

            # Threshold for considering a page as needing OCR
            # If a page has less than 50 characters, it likely needs OCR
            MIN_TEXT_THRESHOLD = 50

            for page_info in pages:
                page_num = page_info.get("page_number", 1)
                page_text = page_info.get("content", "")
                char_count = len(page_text.strip())

                # Check if this page needs OCR
                needs_ocr = (
                    processing_mode == "image_only" or  # Force OCR for all pages
                    (processing_mode == "vlm_enhanced" and char_count < MIN_TEXT_THRESHOLD)
                )

                if needs_ocr and page_num <= total_pages:
                    try:
                        # Convert page to image (2x zoom for OCR)
                        img_bytes = pdf_compat.render_page_png(pdf_content, page_num - 1, zoom=2.0)

                        # Check if page is mostly blank before OCR
                        if pdf_compat.is_blank_page(img_bytes):
                            print(f"[OCR] Page {page_num}: Skipped (blank page)")
                            enhanced_pages.append(page_info)
                            full_text_parts.append(page_text)
                            continue

                        # Call VLM OCR
                        ocr_result = await self.vlm_service.extract_text_ocr(
                            img_bytes, language=language
                        )

                        ocr_text = ocr_result.get("text", "")
                        ocr_confidence = ocr_result.get("confidence", 0)

                        # Validate OCR result to detect hallucination
                        if ocr_text and len(ocr_text) > char_count:
                            if self._is_hallucinated_ocr(ocr_text, language):
                                print(f"[OCR] Page {page_num}: Skipped (hallucination detected)")
                                enhanced_pages.append(page_info)
                                full_text_parts.append(page_text)
                                continue

                            # Use OCR result if it's better than pypdf extraction
                            enhanced_pages.append({
                                "page_number": page_num,
                                "content": ocr_text,
                                "char_count": len(ocr_text),
                                "ocr_applied": True,
                                "ocr_confidence": ocr_confidence,
                                "original_char_count": char_count
                            })
                            full_text_parts.append(ocr_text)
                            ocr_count += 1
                            print(f"[OCR] Page {page_num}: {char_count} → {len(ocr_text)} chars (VLM OCR)")
                        else:
                            # Keep original if OCR didn't improve
                            enhanced_pages.append(page_info)
                            full_text_parts.append(page_text)

                    except Exception as ocr_error:
                        print(f"[OCR] Page {page_num} OCR failed: {ocr_error}")
                        enhanced_pages.append(page_info)
                        full_text_parts.append(page_text)
                else:
                    # Keep original page
                    enhanced_pages.append(page_info)
                    full_text_parts.append(page_text)

            # Update result with enhanced pages
            parsed_result["pages"] = enhanced_pages
            parsed_result["text_content"] = "\n\n".join(full_text_parts)
            parsed_result["metadata"]["ocr_pages_count"] = ocr_count
            parsed_result["metadata"]["vlm_enhanced"] = True

            if ocr_count > 0:
                print(f"[OCR] Enhanced {ocr_count} pages with VLM OCR")

            return parsed_result

        except Exception as e:
            print(f"[OCR] VLM enhancement failed: {e}")
            return parsed_result

    def _is_blank_page(self, png_bytes: bytes) -> bool:
        """
        Check if a page image is mostly blank/white.

        Args:
            png_bytes: PNG image bytes

        Returns:
            True if page is mostly blank (>95% white pixels)
        """
        try:
            from . import pdf_compat
            return pdf_compat.is_blank_page(png_bytes)
        except Exception:
            return False

    # Legacy compatibility - kept for reference but no longer called
    def _is_blank_page_legacy(self, pixmap) -> bool:
        """Legacy PyMuPDF pixmap-based blank page detection (unused)."""
        try:
            samples = pixmap.samples
            width, height = pixmap.width, pixmap.height
            n = pixmap.n

            white_threshold = 250
            total_pixels = width * height
            white_pixels = 0

            step = 10 * n
            for i in range(0, len(samples) - n + 1, step):
                # Check if pixel is white (all components >= threshold)
                is_white = all(samples[i + j] >= white_threshold for j in range(min(n, 3)))
                if is_white:
                    white_pixels += 1

            sampled_total = total_pixels // 10
            white_ratio = white_pixels / sampled_total if sampled_total > 0 else 0

            return white_ratio > 0.95  # More than 95% white = blank page

        except Exception:
            return False  # If check fails, assume not blank

    def _is_hallucinated_ocr(self, ocr_text: str, language: str = "auto") -> bool:
        """
        Detect if OCR result appears to be hallucinated content.

        Args:
            ocr_text: The OCR extracted text
            language: Expected document language

        Returns:
            True if the text appears to be hallucinated
        """
        if not ocr_text:
            return False

        # Hallucination patterns - content unrelated to document context
        hallucination_patterns = [
            # Generic data tables that VLMs often hallucinate
            r"Number of employees",
            r"Characteristic\s*\|\s*Number",
            r"\d{4}\s*\|\s*\d{4}",  # Year | Number pattern
            r"population|revenue|sales|growth",
            # Lorem ipsum and placeholder text
            r"lorem ipsum",
            r"sample text",
            r"placeholder",
            # Generic AI responses
            r"I cannot see|I don't see|image is blank|no text visible",
            r"I'm unable to|I am unable to",
        ]

        text_lower = ocr_text.lower()

        for pattern in hallucination_patterns:
            if re.search(pattern, text_lower, re.IGNORECASE):
                return True

        # Check for suspicious repetitive content
        # If the same phrase repeats too many times, likely hallucination
        words = ocr_text.split()
        if len(words) > 20:
            word_counts = {}
            for word in words:
                word_counts[word] = word_counts.get(word, 0) + 1

            # If any word appears more than 30% of the time, suspicious
            max_count = max(word_counts.values())
            if max_count / len(words) > 0.3:
                return True

        # Language mismatch detection
        if language == "ja":
            # Japanese document should have some Japanese characters
            jp_chars = sum(1 for c in ocr_text if '\u3040' <= c <= '\u30ff' or '\u4e00' <= c <= '\u9fff')
            if len(ocr_text) > 100 and jp_chars < len(ocr_text) * 0.1:
                # Less than 10% Japanese in a Japanese document = suspicious
                return True
        elif language == "ko":
            # Korean document should have some Hangul
            ko_chars = sum(1 for c in ocr_text if '\uac00' <= c <= '\ud7af')
            if len(ocr_text) > 100 and ko_chars < len(ocr_text) * 0.1:
                return True

        return False

    def _sync_extract_pdf(
        self,
        content: bytes,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Synchronous PDF extraction using pypdf."""
        from pypdf import PdfReader

        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)

        # Extract real page count
        page_count = len(reader.pages)

        # Extract text from each page
        pages = []
        full_text = []
        image_count = 0

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            pages.append({
                "page_number": i + 1,
                "content": page_text,
                "char_count": len(page_text)
            })
            full_text.append(page_text)

            # Count images in page (pypdf can detect image objects)
            try:
                if "/XObject" in page["/Resources"]:
                    x_objects = page["/Resources"]["/XObject"].get_object()
                    for obj_name in x_objects:
                        obj = x_objects[obj_name]
                        if obj["/Subtype"] == "/Image":
                            image_count += 1
            except (KeyError, TypeError, AttributeError):
                pass

        # Extract metadata
        pdf_metadata = reader.metadata or {}
        metadata = {
            "pages": page_count,
            "title": pdf_metadata.get("/Title", "").strip() if pdf_metadata.get("/Title") else "",
            "author": pdf_metadata.get("/Author", "").strip() if pdf_metadata.get("/Author") else "",
            "subject": pdf_metadata.get("/Subject", "").strip() if pdf_metadata.get("/Subject") else "",
            "creator": pdf_metadata.get("/Creator", "").strip() if pdf_metadata.get("/Creator") else "",
            "created_date": str(pdf_metadata.get("/CreationDate", "")) if pdf_metadata.get("/CreationDate") else "",
            "modified_date": str(pdf_metadata.get("/ModDate", "")) if pdf_metadata.get("/ModDate") else "",
            "pdf_version": f"{reader.pdf_header}" if hasattr(reader, 'pdf_header') else "unknown"
        }

        # Use filename as title if not available
        if not metadata["title"]:
            metadata["title"] = options.get("filename", "Untitled").replace(".pdf", "")

        # Create image info list with actual image data
        images = []
        if options.get("extract_images", True):
            try:
                from . import pdf_compat

                total_pages = pdf_compat.get_page_count(content)
                img_idx = 0
                for page_num in range(total_pages):
                    image_list = pdf_compat.get_images_metadata(content, page_num)
                    for img_index, img in enumerate(image_list):
                        try:
                            base_image = pdf_compat.extract_image_data(content, page_num, img_index)
                            if base_image:
                                image_bytes = base_image["image"]
                                image_ext = base_image["ext"]
                                width = base_image.get("width", 0)
                                height = base_image.get("height", 0)
                                # Only include images larger than 100x100
                                if width >= 100 and height >= 100:
                                    images.append(ImageInfo(
                                        id=f"img_{uuid.uuid4().hex[:8]}",
                                        page_number=page_num + 1,
                                        position={"x": 0, "y": 0, "width": width, "height": height},
                                        description=f"Image on page {page_num + 1}",
                                        alt_text=f"Embedded image {img_idx + 1}",
                                        data=image_bytes,  # Actual image binary data
                                        width=width,
                                        height=height,
                                        mime_type=f"image/{image_ext}"
                                    ))
                                    img_idx += 1
                        except Exception as img_err:
                            pass  # Skip problematic images
            except Exception as e:
                print(f"Image extraction error: {e}")

        # Extract figure references and match to images
        if images:
            try:
                from .figure_reference_extractor import get_figure_extractor
                extractor = get_figure_extractor()

                # Extract figure references from page text
                figure_refs = extractor.extract_references(pages)

                if figure_refs:
                    # Convert images to dict format for matching
                    image_dicts = [{"page_number": img.page_number} for img in images]

                    # Match references to images
                    mappings = extractor.match_to_images(figure_refs, image_dicts)

                    # Update images with figure references
                    for mapping in mappings:
                        if 0 <= mapping.image_index < len(images):
                            img = images[mapping.image_index]
                            # Update the ImageInfo with figure reference
                            images[mapping.image_index] = ImageInfo(
                                id=img.id,
                                page_number=img.page_number,
                                position=img.position,
                                description=img.description,
                                alt_text=img.alt_text,
                                data=img.data,
                                width=img.width,
                                height=img.height,
                                mime_type=img.mime_type,
                                figure_reference=mapping.normalized,
                                figure_caption=mapping.caption
                            )
                            print(f"Matched figure {mapping.reference} -> image {mapping.image_index} (page {img.page_number})")
            except Exception as e:
                print(f"Figure reference extraction error: {e}")

        # Table detection (basic heuristic - count lines with multiple | or tabs)
        tables = []
        table_pattern = re.compile(r'(\|.*\|)|(\t.*\t.*\t)')
        table_count = 0
        for i, page in enumerate(pages):
            content = page.get("content", "")
            if table_pattern.search(content):
                table_count += 1
                tables.append(TableInfo(
                    id=f"table_{uuid.uuid4().hex[:8]}",
                    page_number=i + 1,
                    headers=[],
                    rows=[],
                    caption=f"Table detected on page {i + 1}",
                    markdown=""
                ))

        return {
            "text_content": "\n\n".join(full_text),
            "pages": pages,
            "metadata": metadata,
            "images": images,
            "tables": tables
        }

    async def _fallback_mock_extraction(self, content: bytes) -> Dict[str, Any]:
        """Fallback mock extraction when pypdf fails."""
        await asyncio.sleep(0.1)
        return {
            "text_content": "PDF content could not be extracted",
            "pages": [{"page_number": 1, "content": "Content extraction failed", "char_count": 0}],
            "metadata": {
                "pages": 1,
                "title": "Unknown",
                "author": "",
                "subject": "",
                "creator": "",
                "created_date": "",
                "modified_date": "",
                "pdf_version": "unknown"
            },
            "images": [],
            "tables": []
        }


class WordParser(BaseDocumentParser):
    """Parser for Microsoft Word documents using python-docx."""

    SUPPORTED_TYPES = [
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ]

    # Heading style patterns (English and Korean)
    HEADING_STYLES = {
        "Heading 1": 1, "Heading 2": 2, "Heading 3": 3,
        "Heading 4": 4, "Heading 5": 5, "Heading 6": 6,
        "제목 1": 1, "제목 2": 2, "제목 3": 3,
        "제목 4": 4, "제목 5": 5, "제목 6": 6,
        "Title": 0, "Subtitle": 1,
    }

    # Chunking configuration
    MAX_SECTION_CHARS = 2000

    def supports(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        """
        Parse Word document using python-docx.

        Extracts:
        - Paragraphs with heading hierarchy
        - Tables with markdown conversion
        - Embedded images (base64 for VLM)
        - Document metadata (author, title, dates)
        """
        options = options or {}
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if filename.endswith(".doc"):
            mime_type = "application/msword"

        doc = ParsedDocument(DocumentType.WORD, filename, mime_type)

        if not DOCX_AVAILABLE:
            doc.text_content = await self._fallback_extraction(file_content, filename)
            doc.metadata = {"title": filename, "error": "python-docx not available"}
        else:
            # Real extraction using python-docx
            extracted = await self._extract_document(file_content, filename, options)
            doc.text_content = extracted["text_content"]
            doc.metadata = extracted["metadata"]
            doc.tables = extracted["tables"]
            doc.images = extracted["images"]
            doc.pages = extracted["sections"]

        doc.processing_info = {
            "parser": "WordParser",
            "parsed_at": datetime.now(timezone.utc).isoformat(),
            "is_docx": filename.endswith(".docx"),
            "docx_available": DOCX_AVAILABLE
        }

        return doc

    async def _extract_document(
        self,
        content: bytes,
        filename: str,
        options: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Extract content from Word document."""
        try:
            docx = DocxDocument(io.BytesIO(content))

            # Extract metadata
            metadata = self._extract_metadata(docx, filename)

            # Extract paragraphs with structure
            sections = self._extract_sections(docx)

            # Build text content with hierarchy
            text_content = self._build_text_content(sections)

            # Extract tables
            tables = []
            if options.get("extract_tables", True):
                tables = self._extract_tables(docx)

            # Extract images
            images = []
            if options.get("extract_images", True):
                images = self._extract_images(docx)

            # Update metadata with counts
            metadata["section_count"] = len([s for s in sections if s.get("is_heading")])
            metadata["table_count"] = len(tables)
            metadata["image_count"] = len(images)

            return {
                "text_content": text_content,
                "metadata": metadata,
                "tables": tables,
                "images": images,
                "sections": sections
            }

        except Exception as e:
            import logging
            logging.error(f"Error parsing Word document {filename}: {e}")
            return {
                "text_content": f"문서 파싱 오류: {str(e)}",
                "metadata": {"title": filename, "error": str(e)},
                "tables": [],
                "images": [],
                "sections": []
            }

    def _extract_metadata(self, docx: 'DocxDocument', filename: str) -> Dict[str, Any]:
        """Extract document metadata from core properties."""
        props = docx.core_properties

        metadata = {
            "title": props.title or filename.replace(".docx", "").replace(".doc", ""),
            "author": props.author or "Unknown",
            "created_date": props.created.isoformat() if props.created else None,
            "modified_date": props.modified.isoformat() if props.modified else None,
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "category": props.category or "",
            "comments": props.comments or "",
        }

        # Count words and paragraphs
        total_text = ""
        para_count = 0
        for para in docx.paragraphs:
            if para.text.strip():
                total_text += para.text + " "
                para_count += 1

        metadata["word_count"] = len(total_text.split())
        metadata["paragraph_count"] = para_count

        return metadata

    def _extract_sections(self, docx: 'DocxDocument') -> List[Dict[str, Any]]:
        """Extract paragraphs organized by sections/headings."""
        sections = []
        current_headings = [""] * 6  # Track heading hierarchy (H1-H6)
        section_idx = 0

        for para in docx.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else "Normal"
            heading_level = self.HEADING_STYLES.get(style_name)

            if heading_level is not None and heading_level > 0:
                # Update heading hierarchy
                current_headings[heading_level - 1] = text
                # Clear lower level headings
                for i in range(heading_level, 6):
                    current_headings[i] = ""

                section_idx += 1
                sections.append({
                    "page_number": section_idx,
                    "content": text,
                    "is_heading": True,
                    "heading_level": heading_level,
                    "heading_path": " > ".join(h for h in current_headings[:heading_level] if h),
                    "text_length": len(text)
                })
            else:
                # Regular paragraph
                # Detect list items
                is_list = self._is_list_paragraph(para, style_name)
                prefix = ""
                if is_list:
                    prefix = self._get_list_prefix(para)

                sections.append({
                    "page_number": section_idx if section_idx > 0 else 1,
                    "content": prefix + text,
                    "is_heading": False,
                    "heading_level": None,
                    "heading_path": " > ".join(h for h in current_headings if h),
                    "is_list": is_list,
                    "text_length": len(text)
                })

        return sections

    def _is_list_paragraph(self, para, style_name: str) -> bool:
        """Check if paragraph is a list item."""
        if "List" in style_name or "Bullet" in style_name or "목록" in style_name:
            return True
        # Check for numbering in XML
        try:
            numPr = para._p.pPr.numPr if para._p.pPr is not None else None
            return numPr is not None
        except Exception:
            return False

    def _get_list_prefix(self, para) -> str:
        """Get appropriate list prefix (bullet or number)."""
        style_name = para.style.name if para.style else ""
        if "Bullet" in style_name or "bullet" in style_name.lower():
            return "• "
        # For numbered lists, we can't easily get the number without complex XML parsing
        # So we use a generic bullet for now
        return "• "

    def _build_text_content(self, sections: List[Dict[str, Any]]) -> str:
        """Build formatted text content from sections."""
        text_parts = []
        current_heading_path = ""

        for section in sections:
            if section.get("is_heading"):
                level = section.get("heading_level", 1)
                prefix = "#" * level
                text_parts.append(f"\n{prefix} {section['content']}\n")
                current_heading_path = section.get("heading_path", "")
            else:
                content = section["content"]
                if section.get("is_list"):
                    text_parts.append(content)
                else:
                    text_parts.append(content)

        return "\n".join(text_parts)

    def _extract_tables(self, docx: 'DocxDocument') -> List[TableInfo]:
        """Extract tables and convert to markdown."""
        tables = []

        for table_idx, table in enumerate(docx.tables):
            rows_data = []

            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    # Handle merged cells by checking for duplicate text
                    cell_text = cell.text.strip().replace("\n", " ")
                    row_cells.append(cell_text)
                rows_data.append(row_cells)

            if not rows_data:
                continue

            # First row as headers
            headers = rows_data[0] if rows_data else []
            data_rows = rows_data[1:] if len(rows_data) > 1 else []

            # Build markdown
            markdown = self._rows_to_markdown(headers, data_rows)

            tables.append(TableInfo(
                id=f"table_{uuid.uuid4().hex[:8]}",
                page_number=table_idx + 1,
                headers=headers,
                rows=data_rows,
                caption=f"테이블 {table_idx + 1}",
                markdown=markdown
            ))

        return tables

    def _rows_to_markdown(self, headers: List[str], rows: List[List[str]]) -> str:
        """Convert table to markdown format."""
        if not headers:
            return ""

        def escape_pipe(s: str) -> str:
            return s.replace("|", "\\|")

        md_lines = []

        # Header
        md_lines.append("| " + " | ".join(escape_pipe(h) for h in headers) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

        # Rows
        for row in rows:
            padded = row + [""] * (len(headers) - len(row)) if len(row) < len(headers) else row[:len(headers)]
            md_lines.append("| " + " | ".join(escape_pipe(str(c)) for c in padded) + " |")

        return "\n".join(md_lines)

    def _extract_images(self, docx: 'DocxDocument') -> List[ImageInfo]:
        """Extract embedded images from document."""
        images = []
        image_idx = 0

        for rel in docx.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type

                    # Get filename/extension
                    filename = getattr(image_part, 'partname', f'image_{image_idx}')
                    if hasattr(filename, 'split'):
                        filename = filename.split('/')[-1]

                    image_idx += 1

                    images.append(ImageInfo(
                        id=f"img_{uuid.uuid4().hex[:8]}",
                        page_number=image_idx,
                        position={},
                        description=f"문서 이미지 {image_idx}",
                        alt_text=str(filename),
                        data=image_bytes,
                        mime_type=content_type or "image/png"
                    ))

                except Exception as e:
                    continue

        return images

    async def _fallback_extraction(self, content: bytes, filename: str) -> str:
        """Fallback when python-docx is not available."""
        return f"Word 문서 '{filename}'을 파싱할 수 없습니다. python-docx 라이브러리가 필요합니다."


class ExcelParser(BaseDocumentParser):
    """Parser for Microsoft Excel documents using openpyxl."""

    SUPPORTED_TYPES = [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ]

    # Chunking configuration
    ROWS_PER_CHUNK = 50  # Maximum rows per chunk for large sheets
    CHUNK_OVERLAP_ROWS = 2  # Overlap rows between chunks for context

    def supports(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        """
        Parse Excel document using openpyxl.

        Extracts:
        - Sheet names and data
        - Headers (first row)
        - Data rows with type conversion
        - Merged cell handling
        - Large sheet chunking with header context
        """
        options = options or {}
        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        if filename.endswith(".xls"):
            mime_type = "application/vnd.ms-excel"

        doc = ParsedDocument(DocumentType.EXCEL, filename, mime_type)

        if not OPENPYXL_AVAILABLE:
            sheets = await self._fallback_extraction(file_content, filename)
        else:
            sheets = await self._extract_sheets(file_content, filename, options)

        # Convert sheets to text and tables
        text_parts = []
        tables = []
        all_chunks = []

        for sheet in sheets:
            sheet_name = sheet["name"]
            headers = sheet["headers"]
            rows = sheet["rows"]
            row_count = len(rows)

            # Sheet header
            text_parts.append(f"=== 시트: {sheet_name} ({row_count}행) ===")

            # Chunk large sheets
            if row_count > self.ROWS_PER_CHUNK:
                chunks = self._chunk_rows(headers, rows, sheet_name)
                for chunk in chunks:
                    table = TableInfo(
                        id=f"table_{uuid.uuid4().hex[:8]}",
                        page_number=sheet.get("sheet_index", 1),
                        headers=headers,
                        rows=chunk["rows"],
                        caption=f"시트: {sheet_name} (행 {chunk['start_row']}-{chunk['end_row']})",
                        markdown=self._rows_to_markdown(headers, chunk["rows"])
                    )
                    tables.append(table)
                    all_chunks.append({
                        "sheet": sheet_name,
                        "start_row": chunk["start_row"],
                        "end_row": chunk["end_row"],
                        "row_count": len(chunk["rows"])
                    })
                    text_parts.append(f"\n[청크 {chunk['start_row']}-{chunk['end_row']}]")
                    text_parts.append(table.markdown)
            else:
                # Small sheet - single table
                table = TableInfo(
                    id=f"table_{uuid.uuid4().hex[:8]}",
                    page_number=sheet.get("sheet_index", 1),
                    headers=headers,
                    rows=rows,
                    caption=f"시트: {sheet_name}",
                    markdown=self._rows_to_markdown(headers, rows)
                )
                tables.append(table)
                text_parts.append(table.markdown)

        doc.text_content = "\n\n".join(text_parts)
        doc.tables = tables

        # Calculate statistics
        total_rows = sum(len(s["rows"]) for s in sheets)
        max_cols = max((len(s["headers"]) for s in sheets), default=0)

        doc.metadata = {
            "title": filename.replace(".xlsx", "").replace(".xls", ""),
            "sheet_count": len(sheets),
            "sheet_names": [s["name"] for s in sheets],
            "total_rows": total_rows,
            "total_columns": max_cols,
            "total_tables": len(tables),
            "chunked": len(all_chunks) > 0,
            "chunk_info": all_chunks if all_chunks else None
        }

        doc.pages = [
            {
                "page_number": i + 1,
                "content": s["name"],
                "row_count": len(s["rows"]),
                "col_count": len(s["headers"])
            }
            for i, s in enumerate(sheets)
        ]

        doc.processing_info = {
            "parser": "ExcelParser",
            "parsed_at": datetime.now(timezone.utc).isoformat(),
            "is_xlsx": filename.endswith(".xlsx"),
            "openpyxl_available": OPENPYXL_AVAILABLE
        }

        return doc

    async def _extract_sheets(
        self,
        content: bytes,
        filename: str,
        options: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """Extract sheets using openpyxl."""
        sheets_data = []

        try:
            # Load workbook with data_only to get calculated values
            wb = load_workbook(
                io.BytesIO(content),
                data_only=True,
                read_only=False  # Need full access for merged cells
            )

            skip_hidden = options.get("skip_hidden_sheets", True)
            skip_empty = options.get("skip_empty_sheets", True)

            for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
                ws = wb[sheet_name]

                # Skip hidden sheets if configured
                if skip_hidden and ws.sheet_state == 'hidden':
                    continue

                # Get dimensions
                if ws.max_row is None or ws.max_row == 0:
                    if skip_empty:
                        continue
                    sheets_data.append({
                        "name": sheet_name,
                        "sheet_index": sheet_idx,
                        "headers": [],
                        "rows": []
                    })
                    continue

                # Extract headers (first row)
                headers = []
                for col in range(1, (ws.max_column or 1) + 1):
                    cell_value = ws.cell(row=1, column=col).value
                    headers.append(self._cell_to_string(cell_value) if cell_value is not None else f"Column{col}")

                # Extract data rows
                rows = []
                for row_idx in range(2, (ws.max_row or 1) + 1):
                    row_data = []
                    is_empty_row = True

                    for col_idx in range(1, (ws.max_column or 1) + 1):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        cell_value = self._get_cell_value(cell, ws)
                        row_data.append(cell_value)
                        if cell_value:
                            is_empty_row = False

                    # Skip completely empty rows
                    if not is_empty_row:
                        rows.append(row_data)

                # Skip empty sheets
                if skip_empty and not rows:
                    continue

                sheets_data.append({
                    "name": sheet_name,
                    "sheet_index": sheet_idx,
                    "headers": headers,
                    "rows": rows,
                    "merged_cells": self._get_merged_ranges(ws)
                })

            wb.close()

        except Exception as e:
            import logging
            logging.error(f"Error parsing Excel {filename}: {e}")
            if not sheets_data:
                sheets_data = await self._fallback_extraction(content, filename)

        return sheets_data

    def _get_cell_value(self, cell, worksheet) -> str:
        """Get cell value, handling merged cells."""
        value = cell.value

        # Check if cell is part of merged range
        for merged_range in worksheet.merged_cells.ranges:
            if cell.coordinate in merged_range:
                # Get value from top-left cell of merged range
                value = worksheet.cell(
                    row=merged_range.min_row,
                    column=merged_range.min_col
                ).value
                break

        return self._cell_to_string(value)

    def _cell_to_string(self, value) -> str:
        """Convert cell value to string."""
        if value is None:
            return ""
        if isinstance(value, bool):
            return "TRUE" if value else "FALSE"
        if isinstance(value, (int, float)):
            # Handle integers and floats
            if isinstance(value, float) and value.is_integer():
                return str(int(value))
            return str(value)
        if hasattr(value, 'strftime'):
            # Handle datetime
            return value.strftime("%Y-%m-%d %H:%M:%S") if hasattr(value, 'hour') else value.strftime("%Y-%m-%d")
        return str(value).strip()

    def _get_merged_ranges(self, worksheet) -> List[str]:
        """Get list of merged cell ranges."""
        return [str(r) for r in worksheet.merged_cells.ranges]

    def _chunk_rows(
        self,
        headers: List[str],
        rows: List[List[str]],
        sheet_name: str
    ) -> List[Dict[str, Any]]:
        """Chunk large sheets into smaller pieces with header context."""
        chunks = []
        total_rows = len(rows)
        chunk_size = self.ROWS_PER_CHUNK
        overlap = self.CHUNK_OVERLAP_ROWS

        start_idx = 0
        chunk_num = 1

        while start_idx < total_rows:
            end_idx = min(start_idx + chunk_size, total_rows)

            chunks.append({
                "chunk_num": chunk_num,
                "start_row": start_idx + 2,  # +2 because row 1 is header, data starts at 2
                "end_row": end_idx + 1,
                "rows": rows[start_idx:end_idx]
            })

            # Move to next chunk with overlap
            start_idx = end_idx - overlap if end_idx < total_rows else total_rows
            chunk_num += 1

        return chunks

    def _rows_to_markdown(self, headers: List[str], rows: List[List[str]]) -> str:
        """Convert table rows to markdown format."""
        if not headers:
            return ""

        # Escape pipe characters in content
        def escape_pipe(s: str) -> str:
            return s.replace("|", "\\|")

        md_lines = []

        # Header row
        header_line = "| " + " | ".join(escape_pipe(h) for h in headers) + " |"
        md_lines.append(header_line)

        # Separator
        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

        # Data rows
        for row in rows:
            # Pad row if needed
            padded_row = row + [""] * (len(headers) - len(row)) if len(row) < len(headers) else row[:len(headers)]
            row_line = "| " + " | ".join(escape_pipe(str(cell)) for cell in padded_row) + " |"
            md_lines.append(row_line)

        return "\n".join(md_lines)

    async def _fallback_extraction(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Fallback extraction when openpyxl is not available."""
        return [
            {
                "name": "Sheet1",
                "sheet_index": 1,
                "headers": ["Error"],
                "rows": [[f"Excel 파일 '{filename}'을 파싱할 수 없습니다. openpyxl 라이브러리가 필요합니다."]]
            }
        ]


class PowerPointParser(BaseDocumentParser):
    """Parser for Microsoft PowerPoint documents using python-pptx."""

    SUPPORTED_TYPES = [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    def supports(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        """
        Parse PowerPoint document using python-pptx.

        Extracts:
        - Slide text (title, body, shapes)
        - Speaker notes
        - Tables
        - Images (as base64 for VLM processing)
        """
        options = options or {}
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if filename.endswith(".ppt"):
            mime_type = "application/vnd.ms-powerpoint"

        doc = ParsedDocument(DocumentType.POWERPOINT, filename, mime_type)

        if not PPTX_AVAILABLE:
            # Fallback to mock if python-pptx not installed
            slides = await self._fallback_extraction(file_content, filename)
        else:
            # Real extraction using python-pptx
            slides = await self._extract_slides(file_content, filename)

        # Convert slides to text
        text_parts = []
        for slide in slides:
            slide_num = slide["number"]
            title = slide.get("title", f"슬라이드 {slide_num}")
            content = slide.get("content", "")

            text_parts.append(f"--- 슬라이드 {slide_num}: {title} ---")
            if content:
                text_parts.append(content)

            # Add speaker notes if available
            if slide.get("notes"):
                text_parts.append(f"[발표자 노트] {slide['notes']}")

            # Add table markdown if present
            if slide.get("tables"):
                for table_md in slide["tables"]:
                    text_parts.append(f"[테이블]\n{table_md}")

        doc.text_content = "\n\n".join(text_parts)

        doc.metadata = {
            "title": filename.replace(".pptx", "").replace(".ppt", ""),
            "slide_count": len(slides),
            "has_notes": any(s.get("notes") for s in slides),
            "has_tables": any(s.get("tables") for s in slides),
            "has_images": any(s.get("images") for s in slides),
            "total_images": sum(len(s.get("images", [])) for s in slides),
            "total_tables": sum(len(s.get("tables", [])) for s in slides)
        }

        doc.pages = [
            {
                "page_number": s["number"],
                "content": s.get("title", ""),
                "text_length": len(s.get("content", ""))
            }
            for s in slides
        ]

        # Extract images from slides
        if options.get("extract_images", True):
            doc.images = await self._collect_images(slides)

        # Extract tables
        doc.tables = await self._collect_tables(slides)

        doc.processing_info = {
            "parser": "PowerPointParser",
            "parsed_at": datetime.now(timezone.utc).isoformat(),
            "is_pptx": filename.endswith(".pptx"),
            "pptx_available": PPTX_AVAILABLE
        }

        return doc

    async def _extract_slides(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Extract slides using python-pptx."""
        slides_data = []

        try:
            prs = Presentation(io.BytesIO(content))

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_info = {
                    "number": slide_idx,
                    "title": "",
                    "content": "",
                    "notes": "",
                    "tables": [],
                    "images": []
                }

                # Extract title
                if slide.shapes.title:
                    slide_info["title"] = slide.shapes.title.text.strip()

                # Extract text from all shapes
                text_parts = []
                for shape in slide.shapes:
                    # Skip title shape (already extracted)
                    if shape == slide.shapes.title:
                        continue

                    # Extract text from text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = "".join(run.text for run in paragraph.runs).strip()
                            if para_text:
                                # Detect bullet points
                                if paragraph.level > 0:
                                    para_text = "  " * paragraph.level + "• " + para_text
                                elif self._is_bullet_paragraph(paragraph):
                                    para_text = "• " + para_text
                                text_parts.append(para_text)

                    # Extract tables
                    if shape.has_table:
                        table_md = self._extract_table(shape.table)
                        slide_info["tables"].append(table_md)

                    # Extract images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_info = self._extract_image(shape, slide_idx)
                        if img_info:
                            slide_info["images"].append(img_info)

                slide_info["content"] = "\n".join(text_parts)

                # Extract speaker notes
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame:
                        notes_text = notes_frame.text.strip()
                        if notes_text:
                            slide_info["notes"] = notes_text

                slides_data.append(slide_info)

        except Exception as e:
            # Log error and return empty or partial data
            import logging
            logging.error(f"Error parsing PowerPoint {filename}: {e}")
            if not slides_data:
                slides_data = await self._fallback_extraction(content, filename)

        return slides_data

    def _is_bullet_paragraph(self, paragraph) -> bool:
        """Check if paragraph has bullet formatting."""
        try:
            # Check for bullet character in XML
            pPr = paragraph._p.get_or_add_pPr()
            buNone = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
            return buNone is None and paragraph.level >= 0
        except Exception:
            return False

    def _extract_table(self, table) -> str:
        """Extract table and convert to markdown format."""
        rows = []
        for row_idx, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            rows.append(cells)

        if not rows:
            return ""

        # Build markdown table
        md_lines = []

        # Header row
        md_lines.append("| " + " | ".join(rows[0]) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")

        # Data rows
        for row in rows[1:]:
            # Pad row if needed
            while len(row) < len(rows[0]):
                row.append("")
            md_lines.append("| " + " | ".join(row[:len(rows[0])]) + " |")

        return "\n".join(md_lines)

    def _extract_image(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract image from shape."""
        try:
            image = shape.image
            image_bytes = image.blob
            content_type = image.content_type

            # Get position and size
            position = {
                "x": shape.left,
                "y": shape.top,
                "width": shape.width,
                "height": shape.height
            }

            # Convert to base64 for VLM processing
            image_b64 = base64.b64encode(image_bytes).decode('utf-8')

            return {
                "id": f"img_{uuid.uuid4().hex[:8]}",
                "slide_number": slide_num,
                "content_type": content_type,
                "position": position,
                "base64": image_b64,
                "size_bytes": len(image_bytes)
            }
        except Exception as e:
            return None

    async def _collect_images(self, slides: List[Dict[str, Any]]) -> List[ImageInfo]:
        """Collect all images from slides into ImageInfo list."""
        images = []
        for slide in slides:
            for img in slide.get("images", []):
                # Decode base64 to bytes for the data field
                img_data = None
                if img.get("base64"):
                    try:
                        img_data = base64.b64decode(img["base64"])
                    except Exception:
                        pass

                images.append(ImageInfo(
                    id=img["id"],
                    page_number=img["slide_number"],
                    position=img.get("position", {}),
                    description=f"슬라이드 {img['slide_number']}의 이미지",
                    alt_text=f"Image from slide {img['slide_number']}",
                    data=img_data,
                    mime_type=img.get("content_type", "image/png")
                ))
        return images

    async def _collect_tables(self, slides: List[Dict[str, Any]]) -> List[TableInfo]:
        """Collect all tables from slides into TableInfo list."""
        tables = []
        for slide in slides:
            for table_idx, table_md in enumerate(slide.get("tables", [])):
                # Parse markdown to extract headers and rows
                lines = table_md.strip().split("\n")
                headers = []
                rows = []

                if lines:
                    header_line = lines[0].strip()
                    if header_line.startswith("|"):
                        headers = [h.strip() for h in header_line.split("|")[1:-1]]

                    # Extract data rows (skip header and separator line)
                    for line in lines[2:]:
                        if line.strip().startswith("|"):
                            row_cells = [c.strip() for c in line.split("|")[1:-1]]
                            rows.append(row_cells)

                tables.append(TableInfo(
                    id=f"table_{uuid.uuid4().hex[:8]}",
                    page_number=slide["number"],
                    headers=headers,
                    rows=rows,
                    markdown=table_md
                ))
        return tables

    async def _fallback_extraction(self, content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Fallback extraction when python-pptx is not available or parsing fails."""
        return [
            {
                "number": 1,
                "title": "파싱 실패",
                "content": f"PowerPoint 파일 '{filename}'을 파싱할 수 없습니다. python-pptx 라이브러리가 필요합니다.",
                "notes": "",
                "tables": [],
                "images": []
            }
        ]


class ImageParser(BaseDocumentParser):
    """Parser for image files with VLM support."""

    SUPPORTED_TYPES = [
        "image/png", "image/jpeg", "image/jpg", "image/gif",
        "image/bmp", "image/tiff", "image/webp"
    ]

    def __init__(self, vlm_service=None):
        self.vlm_service = vlm_service

    def supports(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        options: Dict[str, Any] = None
    ) -> ParsedDocument:
        """
        Parse image file using VLM for understanding.
        """
        options = options or {}
        ext = os.path.splitext(filename)[1].lower()
        mime_type = EXTENSION_TO_MIME.get(ext, "image/png")

        doc = ParsedDocument(DocumentType.IMAGE, filename, mime_type)

        # If VLM service is available, use it for extraction
        if self.vlm_service and options.get("use_vlm", True):
            vlm_result = await self.vlm_service.process_image(
                file_content,
                task="describe"
            )
            doc.text_content = vlm_result.text_content

            # Extract tables if detected
            if vlm_result.tables:
                for t in vlm_result.tables:
                    doc.tables.append(TableInfo(
                        id=t.get("id", f"table_{uuid.uuid4().hex[:8]}"),
                        headers=t.get("headers", []),
                        rows=t.get("rows", []),
                        caption=t.get("caption", "")
                    ))

            doc.metadata = {
                "vlm_processed": True,
                "confidence": vlm_result.confidence_score,
                "processing_time_ms": vlm_result.processing_time_ms,
                "layout": vlm_result.layout_analysis,
                "detected_objects": vlm_result.detected_objects
            }
        else:
            # Basic image info without VLM
            doc.text_content = f"이미지 파일: {filename}"
            doc.metadata = {
                "vlm_processed": False,
                "format": ext.replace(".", "").upper()
            }

        # Create image info
        doc.images = [
            ImageInfo(
                id=f"img_{uuid.uuid4().hex[:8]}",
                description=doc.text_content[:200] if doc.text_content else filename,
                alt_text=filename
            )
        ]

        doc.pages = [{"page_number": 1, "content": doc.text_content}]

        doc.processing_info = {
            "parser": "ImageParser",
            "parsed_at": datetime.now(timezone.utc).isoformat(),
            "vlm_enabled": bool(self.vlm_service) and options.get("use_vlm", True)
        }

        return doc


class DocumentParserFactory:
    """Factory for creating appropriate document parsers."""

    def __init__(self, vlm_service=None):
        self.vlm_service = vlm_service
        self.parsers: List[BaseDocumentParser] = [
            TextParser(),
            PDFParser(vlm_service),
            WordParser(),
            ExcelParser(),
            PowerPointParser(),
            ImageParser(vlm_service)
        ]

    def get_parser(self, mime_type: str) -> Optional[BaseDocumentParser]:
        """Get appropriate parser for MIME type."""
        for parser in self.parsers:
            if parser.supports(mime_type):
                return parser
        return None

    def get_parser_for_file(self, filename: str) -> Optional[BaseDocumentParser]:
        """Get appropriate parser based on filename extension."""
        ext = os.path.splitext(filename)[1].lower()
        mime_type = EXTENSION_TO_MIME.get(ext)
        if mime_type:
            return self.get_parser(mime_type)
        return None

    async def parse_document(
        self,
        file_content: bytes,
        filename: str,
        mime_type: str = None,
        options: Dict[str, Any] = None
    ) -> Optional[ParsedDocument]:
        """
        Parse a document using the appropriate parser.

        Args:
            file_content: Raw file bytes
            filename: Original filename
            mime_type: MIME type (optional, will be inferred from filename)
            options: Parsing options

        Returns:
            ParsedDocument or None if no parser available
        """
        # Determine MIME type
        if not mime_type:
            ext = os.path.splitext(filename)[1].lower()
            mime_type = EXTENSION_TO_MIME.get(ext)

        if not mime_type:
            return None

        # Get parser
        parser = self.get_parser(mime_type)
        if not parser:
            return None

        # Parse document
        return await parser.parse(file_content, filename, options)

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(EXTENSION_TO_MIME.keys())

    def get_supported_mime_types(self) -> List[str]:
        """Get list of supported MIME types."""
        return list(SUPPORTED_MIME_TYPES.keys())


# Factory function
def get_document_parser_factory(vlm_service=None) -> DocumentParserFactory:
    """Get document parser factory instance."""
    return DocumentParserFactory(vlm_service)
